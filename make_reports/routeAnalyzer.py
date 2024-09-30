import sys,os
import math
import subprocess
from rdkit import Chem
from rdkit import DataStructs
from rdkit.Chem.Fingerprints import FingerprintMols
from rdkit.Chem import Draw, rdMolDescriptors
from rdkit.Chem.Draw import rdMolDraw2D
from IPython.display import SVG

from reportlab.graphics import renderPDF
from reportlab.pdfgen import canvas
from reportlab.pdfbase.pdfmetrics import stringWidth
from reportlab.lib.pagesizes import A4, portrait, landscape
from svglib.svglib import svg2rlg
from reportlab.lib.units import mm
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from PIL import Image
import shutil
import Levenshtein

_web=False
_web=True

if _web:
    dr='/var/www/html/public/images/tmp' # working directory
else:
    dr='tmp' # working directory

_magic_x=650

class openReport:
    def __init__(self,file,scale):
        ext=file.split(".")[1]
        if ext=='pdf':
            print("open pdf:",file)
            self.type=ext
            self.scale=scale
            self.page=canvas.Canvas(file,pagesize=landscape(A4))
            self.font_name="Times-Bold"
            self.height=210*mm
            self.width=297*mm
            self.font_size=4*mm
            self._font_size=4*mm
        else:
            self.type='ppt'
            self.scale=scale
            self.page=Presentation()
            self.page.slide_height=Inches(8.27)
            self.page.slide_width =Inches(11.69)
            self.height=self.page.slide_height
            self.sc=Inches(11.69)/830
#            self.sc=Inches(8.0)/800
            blank_slide_layout = self.page.slide_layouts[6]
            self.slide = self.page.slides.add_slide(blank_slide_layout)
            self.file=file
            self.font_name="Times-Bold"
            self.font_size=14 # 14<-16 
            self._font_size=14 #

    def getSpan(self):
#       scale=_magic_x/(span)
        return _magic_x/self.scale
     
    def setFont(self,font_name=None,font_size=None):# size is relative value to default
        if self.type=='pdf':
            if font_name:
                self.font_name=font_name
            if font_size:
                self.font_size=self._font_size*font_size
            self.page.setFont(self.font_name,self.font_size)
        else:
            if font_size:
                self.font_size=self._font_size*font_size

    def stroke(self):
        if self.type=='pdf':
            self.page.showPage()
        else:
            blank_slide_layout = self.page.slide_layouts[6]
            self.slide = self.page.slides.add_slide(blank_slide_layout)

    def drawImage(self,fn,x,y,scale,bt=False):
        if bt:
            y=y+self.scale*0.55
        else:
            y=y-scale*0.4
#            y=y

        if self.type=='pdf':
            drawing=svg2rlg(fn)
            drawing.renderScale=scale/drawing.width
            renderPDF.draw(drawing, canvas=self.page, x=x, y=y)
        else:
            work=dr+"/komai.svg"

            key="2.0px"
            kkk="4.0px"

            with open(fn,"r") as f:
                all=f.readlines()

            with open(work,"w") as f:
                for line in all:
                    o=line.replace(key,kkk)
                    f.write(o)

            gn=fn.split(".svg")[0]+".jpg"
            if _web:
                subprocess.run(['/usr/bin/convert',work,gn])
            else:
                subprocess.run(['convert',work,gn])
            x0,y0=x*self.sc,self.page.slide_height-(y+scale)*self.sc
        
            self.slide.shapes.add_picture(
                gn,left=x0,top=y0,
#                width=scale*self.sc, height=scale*self.sc
                width=scale*self.sc*0.9, height=scale*self.sc*0.9
    )


    def drawLine(self,pt):
        if self.type=='pdf':
            self.page.line(pt[0],pt[1],pt[2],pt[3])
        else:
            x1,y1=pt[0]*self.sc,self.page.slide_height-pt[1]*self.sc
            x2,y2=pt[2]*self.sc,self.page.slide_height-pt[3]*self.sc
            putLine(x1,y1,x2,y2,self.slide.shapes)

    def drawArrow(self,pt):
        if self.type=='pdf':
            self.page.line(pt[0],pt[1],pt[2],pt[3])
            self.page.line(pt[2],pt[3],pt[4],pt[5])
            self.page.line(pt[2],pt[3],pt[6],pt[7])
        else:
            x0,y0=pt[0]*self.sc,self.page.slide_height-pt[1]*self.sc
            x1,y1=pt[2]*self.sc,self.page.slide_height-pt[3]*self.sc
            x2,y2=pt[4]*self.sc,self.page.slide_height-pt[5]*self.sc
            x3,y3=pt[6]*self.sc,self.page.slide_height-pt[7]*self.sc
            group=[]
            group+=[putLine(x0,y0,x1,y1,self.slide.shapes)]
            group+=[putLine(x1,y1,x2,y2,self.slide.shapes)]
            group+=[putLine(x1,y1,x3,y3,self.slide.shapes)]
            self.slide.shapes.add_group_shape(group)

    def getProperLength(self,mag,string):
        if self.type=='pdf':
            tx=string
            while True:
                t=(self.getSpan()-1)*self.scale-stringWidth(tx,self.font_name,self.font_size)
                if t>0:
                    return len(tx)
                tx=tx[0:len(tx)-3]
        else:
            return 50

    def drawString(self,x,y,string,center=False,adj=True):
#        komai
        if self.type=='pdf':
            if center:
                x+=self.getSpan()*self.scale/2-stringWidth(string,self.font_name,self.font_size)/2
            self.page.drawString(x,y,string)
        else:
            xx=x*self.sc
            yy=self.page.slide_height-y*self.sc-self.font_size
            if adj:
#                komai
#                shape=self.slide.shapes.add_textbox(xx,yy-Pt(self.font_size),pw,Cm(1))
                pw=Cm(0.1)*self.font_size*(len(string)+2)*0.18
                py=Cm(0.1)*self.font_size*0.5
                shape=self.slide.shapes.add_textbox(xx,yy-Pt(self.font_size),pw,py)
            else:
                shape=self.slide.shapes.add_textbox(xx,yy-Pt(self.font_size),Cm(1),Cm(10))

            tf=shape.text_frame
#            tf.auto_size=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            p=tf.paragraphs[0]
            run=p.add_run()
            run.text=string
            run.font.name=self.font_name
            run.font.size=Pt(self.font_size)
            #komai

    def close(self):
        if self.type=='pdf':
            self.page.save()
        else:
            self.page.save(self.file)

def putLine(x0,y0,x1,y1,shapes):
    if x1>x0:
        if y1>y0:
           shape=shapes.add_shape(MSO_SHAPE.LINE_INVERSE,x0,y0,x1-x0,y1-y0)
           shape.element.getchildren()[1].getchildren()[0].set('flipV','1')
        else:
           shape=shapes.add_shape(MSO_SHAPE.LINE_INVERSE,x0,y1,x1-x0,-y1+y0)
    else:
        if y1>y0:
           shape=shapes.add_shape(MSO_SHAPE.LINE_INVERSE,x1,y0,-x1+x0,y1-y0)
        else:
           shape=shapes.add_shape(MSO_SHAPE.LINE_INVERSE,x1,y1,-x1+x0,-y1+y0)
           shape.element.getchildren()[1].getchildren()[0].set('flipV','1')
    shape.fill.solid()
    shape.fill.fore_color.rgb=RGBColor(0,255,0)
    shape.line.color.rgb=RGBColor(0,0,0)
    return shape

def picSize():
    return 1.2

def smile2sgv(smiles,name,dr=dr,size=200):
#    if os.path.exists(dr):
#        shutil.rmtree(dr)
    if not os.path.exists(dr):
        os.makedirs(dr)

    mol=Chem.MolFromSmiles(smiles)
#    mw=rdMolDescriptors.CalcExactMolWt(mol)
    mw=mol.GetNumAtoms()
#    mw=Chem.AddHs(mol).GetNumAtoms()
#    fc=int(math.sqrt(mw/5))
    fc=int(math.sqrt(mw)/1.5)
#    print(smiles,mw,fc)
#    fc=int(math.pow(mw/20,0.333))
#    print(mw,fc,name)
    if fc<1:
        fc=1
    view = rdMolDraw2D.MolDraw2DSVG(size*fc,size*fc)

    option=view.drawOptions()
    option.circleAtoms=False
    option.continuousHightlight=False

    tm = rdMolDraw2D.PrepareMolForDrawing(mol)
    view.DrawMolecule(tm)
    view.FinishDrawing()

    svg=view.GetDrawingText()

    with open(dr+"/"+name,'w') as f:
        f.write(svg)

    return fc

def getSmiles(tg,lines):
    key='Route:'+str(tg)
    on=False
    ret=[]
    for l in lines:
        if key in l:
            on=True
            continue
        if on:
            if 'Route' in l:
                if len(ret)<1:
#                    print(l1)
                    ret.append([l1.split(',')[0]])
                return ret
            if not '>' in l:
                l1=l
                continue
            each=l.split(',')
            ret.append(each[:-1])
    return ret

def Debug(c,head=None,sp=None):
    return
    for n,i in enumerate(c):
        if head:
            print(head,end=":")
        if sp:
            print(n,i[sp],end="<--")
        print(n,i)
#    exit()


def chkConnect(n,connect):
    if n<0: 
        return n
    n=connect[n][3]
    while connect[n][1][0]==-1:# ext step: no smiles
        n=connect[n][3]
    return n

def chkJoin(c,s):
    if len(c[2])<2:
        return False
    c[2].sort()
    p1=c[2][-1]
#    print("in checkJoin",p1,s[p1][0],s[p1][1],"<---")
#    print(c[2])
    if s[p1][0]=='':
        return False
    return True

def addCatal(s,p):
    pp=[]
    for m,ss in enumerate(s):
        if ss=='>':
            return pp
        if not m in p:
            pp.append(m)
    return pp

def addCatal2(s,p):
    pp=[]
    for m,ss in enumerate(s):
        if ss=='>':
            return pp
        if not ss in p:
            pp.append(m)
    return pp

def lenForCat(cat,fc):
    sp=0
    for i in cat:
        sp+=fc[i]
    return sp
    
def setBranchLength(connect,goal,fc):
    branch=[]
    right=0

    for n in goal:
        px=0
        while True:
            if not fc:
                sp=len(connect[n][-1])
                connect[n][4]=px=px-(sp+2)
            else:
                if n<len(fc):
                    sp=lenForCat(connect[n][-1],fc[n])+0.5
                    wd=fc[n][connect[n][1][0]]+1
                else:
                    wd=fc[-1]['cheat']
                    sp=2
#   komai
                # for starting material
                #
                connect[n][4]=px=px-(sp+wd)

            if len(connect[n][2])<1:
                if len(branch)<1:
                    break
                else:
                    ww=branch.pop(-1)
                    n,px=ww[0],ww[1]
            elif len(connect[n][2])>1:
                for i in connect[n][2][1:]:
                    branch.append([i,px])
                n=connect[n][2][0]
            else:
                n=connect[n][2][0]
            if right>sp:
                right=sp
    for c in connect:
        c[4]-=right

def initConnect(smiles,fc=False,mode=1):
    connect=[]
    drop={}
    for n,proc in enumerate(smiles):
        e=False
#        print("proc")
        for m, comp in enumerate(proc):
            if ">" in comp:
                e=True
                continue
            if e:
# connect=[smiles of output product,[#ID,], [connected from], connect to, x, y]
                connect.append([comp,[m],[],-1,-1,0,0])
#
# up to here: make template for each steps (arrow)
#
    Debug(connect,head='A')

    for n, prev in enumerate(smiles):
        for m, proc in enumerate(smiles):
            if n !=m and connect[n][0] in proc:
#                print("n",n,m)
                connect[n][3]=m
                connect[m][2].append(n)
#
# up to here: find connection (next arrow)
#
    Debug(connect,head='B')
    goal=[]
    l=len(smiles)
    for i in range(l):
        if connect[i][3]==-1:
            goal.append(i)

    if len(goal)>1:
        print("warning multiple goal")

###################################
    for n,c in enumerate(connect):
        while chkJoin(c,connect):
            p1=c[2].pop(-1)
            drop[n]=connect[p1][0]
            p2=c[3]
            p3=len(connect)
            connect[p1][3]=p3
            connect.append(['',[-1],[p1],p2,-1,0,0])
            if p2<0:
                goal.append(p3)
            else:
                connect[p2][2].append(p3)
#
# up to here: make extra stepping stone
#
# p1 -> [,], p2=To, p3=last(inserted)
# in case of p2=-1 --> trouble
#
    Debug(connect,head='C')
#
# up to here: find head : open edge
#
    start=[]
    branch=[]
    connect[n][4]=connect[n][5]=0;
    depth=1
    for n in goal:
        depth-=1
        while True:
            connect[n][5]=depth
            if len(connect[n][2])<1:
                start.append(n)
                if len(branch)<1:
                    break
                else:
                    m=branch.pop(-1)
                    depth=depth-1
                    #n=connect[m][3]
            else:
                m=connect[n][2][0]
                if len(connect[n][2])>1:
                    branch=branch+connect[n][2][1:]
#        print(n,"->",m)
            n=m
#
# up to here: define tree structure
#
    for n in range(len(connect)):
        connect[n].append([])
    Debug(connect,head='D')

#    Debug(connect,head='A')
#    print("keys",drop)
#    print("keys",list(drop.keys()))
# connect=['smiles',[#ID,], [from], to, x, y, start_chem, [catalyser]]
    for n,s in enumerate(smiles):
        c=connect[n]
        p=[]
        if len(c[2])<1:
            init=whichIsFirstMaterial(s,mode)
            c[6]=init
            p.append(init)# will be modified
#            p.append(0)# will be modified
            pp=addCatal(s,p)
            c[7]+=pp
#            print("upper",n)
        else:#
            d=[]
#            print("lower",n)
            for f in c[2]:#
                d.append(connect[f][0])
            if n in list(drop.keys()):
                d.append(drop[n])
            pp=addCatal2(s,d)
            c[7]+=pp
    Debug(connect,head='E')
#
# up to here: register (catlizers) #id 
#
    if fc:
        if len(connect[i][2])<1:
            fc[-1]['cheat']=0
        else:
            i1=connect[i][2][0]
            i2=connect[i1][1][0]
            fc[-1]['cheat']=fc[i1][i2]

    setBranchLength(connect,goal,fc)
#
# inverse ->> n: upstream, m: downstream (goal to sta]rts)
#
#            connect[m][4]=connect[n][4]-(sp+2)
    Debug(connect,head='F')

    return connect,start
#        smile2sgv(comp,"e"+str(n+1)+"x"+str(m+1)+".svg")

def makeSgv(smiles,type=0):
    size=100
    fcs=[]
    for n,proc in enumerate(smiles):
        fc={}
        for m,comp in enumerate(proc):
            if ">" in comp:
                continue
#            print(comp,"e"+str(n+1)+"x"+str(m+1)+".svg") #debuf
            fc[m]=smile2sgv(comp,"e"+str(n+1)+"x"+str(m+1)+".svg",size=size)
        fcs.append(fc)

    return fcs

def arrow(x0,y0,x1,y1,scale,skip=1.0,fc=False,ugly_patch=True):

    if not fc:
        return prevArrow(x0,y0,x1,y1,scale,ugly_patch=ugly_patch)
#
# ugly patch
#
    if ugly_patch and x0>x1:
        x1=x0+scale*3
        y1=y0
# end
# anyway
    l=math.sqrt((x1-x0)*(x1-x0)+(y1-y0)*(y1-y0))
    rt=(l-scale*skip)/l
    sx=10
    pi=3.1415/12.0 # 30 deg

#from: shift full
    x2=x0+skip*scale
    y2=y0+scale/2
#to: shift full 
    x3=x1
    y3=y1+scale/2

    a=math.atan2(y3-y2,x3-x2)

    xx=x3-sx*math.cos(a+pi)
    yy=y3-sx*math.sin(a+pi)
    ret=[x2,y2,x3,y3,xx,yy]

    xx=x3-sx*math.cos(a-pi)
    yy=y3-sx*math.sin(a-pi)

    return ret+[xx,yy]

def prevArrow(x0,y0,x1,y1,scale,skip=1.0,ugly_patch=True):
#
# ugly patch
#
    if ugly_patch and x0>x1:
        x1=x0+scale*3
        y1=y0
# end
# anyway
    l=math.sqrt((x1-x0)*(x1-x0)+(y1-y0)*(y1-y0))
    rt=(l-scale*skip)/l
    sx=10
    pi=3.1415/12.0 # 30 deg

#from: shift full
    x2=x1+(x0-x1)*(rt-0.10)
    y2=y1+(y0-y1)*(rt-0.10)+scale/2
#to: shift full 
    x3=x0+(x1-x0)*(1.0-0.10)
    y3=y0+(y1-y0)*(1.0-0.10)+scale/2

    a=math.atan2(y1-y0,x1-x0)

    xx=x3-sx*math.cos(a+pi)
    yy=y3-sx*math.sin(a+pi)
    ret=[x2,y2,x3,y3,xx,yy]

    xx=x3-sx*math.cos(a-pi)
    yy=y3-sx*math.sin(a-pi)

    return ret+[xx,yy]

def getName(a,b,dr=dr):
    fn=dr+"/e"+str(a+1)+"x"+str(b)+".svg"#<-- need modification
    return(fn)

def getGoal(smiles,connect):
    goal=[]
    l=len(smiles)
    for i in range(l):
        if connect[i][3]==-1:
            goal.append(i)
    return goal

def getSharpHead(h,c):
    sharpHead=[]

    for n in h:
        branch=[]
        while True:
            if len(c[n][2])<1:#i.e. starting point
                if len(branch)<1:# and no reserved entry
                    break
                else:
                    m=branch.pop(-1) # pop and restart from reserved entry
                    n=c[m][3]
            elif len(c[n][2])>1:# branching 
                branch=branch+c[n][2][1:] # thus reserve others
                m=c[n][2][0] # go forward
            else:
                m=c[n][2][0] # 
            if m>-1 and c[n][4]<c[m][4]:
                sharpHead.append([n,m])
            n=m
            if n<0:# anyway no
                break
    return sharpHead

def rightEdge(i,fc,cnt,l_fc):
    if i[1][0]<0 or cnt>=l_fc:
        return i[4]+len(i[7])+2
    return i[4]+len(i[7])+2+fc[cnt][i[1][0]]

def align(c,smiles,s,fc,tite):
    connect=c
    m=0
    l_fc=len(fc)
    annealed=[]
    for i in c:
        if m>i[4]:
            m=i[4]
    for i in c:
        i[4]-=m;

    loop=True
    depth=0# cleaned floor depth
    while loop:
        d=0
        y0=100
        y1=-y0
        x0=s*2
        loop=False
        ccc=0
        for cnt,i in enumerate(c):
            xx0=rightEdge(i,fc,cnt,l_fc)
            if xx0 > s:# arrow starting point!
                loop=True
                if x0 > i[4]:
                    x0=i[4]
                    ccc=cnt
            elif i[5]<=depth:
                if y1 < i[5]:
                    y1=i[5] # smaller 
                if y0 > i[5]:
                    y0=i[5] # bigger
        if loop:
            for cnt,i in enumerate(c):
                xx0=rightEdge(i,fc,cnt,l_fc)
                if xx0 > s:
                    i[4] -= x0
                    i[5] += (y0-y1-1)
        depth+=(y0-y1-1)

        sharpHead=getSharpHead(getGoal(smiles,c),connect)
        annealHead(sharpHead,connect)

        if tite:
            span=6
        else:
            span=4
        deep=[0]
        for n,vv in enumerate(fc):
            for v in vv.values():
                i=-c[n][5]
                while i>=len(deep):
#                if i>=len(deep):
                    deep.append(0)
                if deep[i]<v:
                    deep[i]=v
        dx=[0]
        for n in range(len(deep)-1):
            dx.append((deep[n]+deep[n+1])/span+dx[n])

        for i in c:
            if -i[5]<len(dx):
                i[5]=-int(dx[-i[5]])
            else:
                i[5]=i[4]
        
def prevEvalHeight(connect):
    depth=0
    for i in range(len(connect)):
        if depth > connect[i][5]:
            depth=connect[i][5]
    return depth

def evalHeight(c,fc):
    if not fc:
        return prevEvalHeight(c)

    depth=0
    for n,vv in enumerate(fc):
        for v in vv.values():
            w=c[n][5]-v/2
            if depth > w:
                depth=w
    return depth

def compSmiles(s1,s2):
    for a,b in zip(s1,s2):
        if a!=b:
            return False
    return True

def annealHead(sh,c):
    for p in sh:
        m=len(c)
        keep=c[p[1]][0]
        #komai 0830
#        c[p[1]][3]=-2 # right edge out: no connection
        c[p[1]][3]=-2 # right edge out: no connection
        c[p[0]][2]=[-1] # right edge in

def _ccn(i):
    if i<10:
        return 1
    if i<100:
        return 2
    return 3

def getPropString(m):
    if len(m)<2:
        s=str(m[0])
#        w=_ccn(m[0])+2
        w=_ccn(m[0])+2
        return s,w

#    s=str(m[-1])+"-"+str(m[0])
    s=str(m[0])+"-"+str(m[-1])
#    w=_ccn(m[0])+_ccn(m[-1])+2
    w=_ccn(m[0])+_ccn(m[-1])+2
    return s,w

def preDraw(init,proc,branch):

    ixx=0
    yy=1
    pos={}
    deep={}

    for i in proc[init]:
        pos[i]=[ixx,yy]
        s,w=getPropString(branch[i])
        ixx=ixx+w
    deep[init]=yy

    right=ixx
    theLeft=0
    y1=yy

    for j in reversed(proc):
        if j==init:
            continue
        deep[j]=0
        first=True
        prev=True
        curr=True
        ixx=-1

        for i in reversed(proc[j]):
            if i in pos:
                ixx=pos[i][0]
                yy=pos[i][1]
                curr=True
            else:
#                if ixx<0:
#                    ixx=right
                if ixx<0:
                    ixx=theLeft
                curr=False
                s,w=getPropString(branch[i])
                if first:
                    y1=y1+1
                    first=False
                yy=y1

                pos[i]=[ixx-w,yy]

            ixx=pos[i][0]
            if ixx<theLeft:
                theLeft=ixx
            prev=curr

            if deep[j]<yy:
                deep[j]=yy
    if theLeft<0:
        for i in pos:
            pos[i][0]-=theLeft

    return pos,deep

def similarityOnRoute(init,proc,branch,all):

    pos,deep=preDraw(init,proc,branch)
    name={}
    for j in proc:
        if deep[j] in name:
            name[deep[j]]+=","+str(j)
        else:
            name[deep[j]]=str(j)

    sim=[]
    fps=False;
    for i in name:
        for j in name[i].split(','):
            tg_route=j
            ssMax=0
            ss=getSmiles(j,all)
            sl=len(ss)
            if sl>ssMax:
                ssMax,tg_route=sl,j
        fpMin=10
        if j!=tg_route:
            ss=getSmiles(tg_route,all)
        if not fps:
            fps=FingerprintMols.FingerprintMol(Chem.MolFromSmiles(ss[-1][-1]))
        else:
            tg=whichIsFirstMaterial(ss[0],_type)
            fpp=FingerprintMols.FingerprintMol(Chem.MolFromSmiles(ss[0][tg]))
            sm=DataStructs.FingerprintSimilarity(fps,fpp)
#            sim.append([tg_route,sm])
#    return sim easy way
# question....
            for smile in ss:
                fpp=FingerprintMols.FingerprintMol(Chem.MolFromSmiles(smile[-1]))
                ss=DataStructs.FingerprintSimilarity(fps,fpp)
                if sm>ss:
                    sm=ss
            sim.append([tg_route,sm])
    print("sim",sim)
    return sim

def head_page(head,page,hope):
#head=[chemical,getSmiles(route[0],all),ox,oy,["statistics will be ..."]]
    ox=40
    oy=480
    mag=3

    page.setFont(font_size=mag)
    page.drawString(ox,oy,head[0])
    scale=page.scale
#
    smiles=head[1][-1][-1]
    fn="head.svg"
    x0=ox
    py=oy-scale*10.5
#    print("smiles --->",smiles)
    smile2sgv(smiles,fn,size=100)
    page.drawImage(dr+"/"+fn,x0,py,scale*10)

    pt=0;l=len(smiles)
#    py=py-mag*page.scale
    mg=2
    page.setFont(font_size=mg)
    wd=page.getProperLength(mg,smiles)
#    wd=7
    tl=smiles
    n=int(math.ceil(l/wd))
    py=(n+0.5)*page.font_size
    while pt<l:
        end=min(pt+wd,l)
        tl=smiles[pt:end]
        py-=page.font_size
        page.drawString(ox,py,tl,center=True)
        pt+=wd

#    print("head",type(head[2][0]))
    if type(head[2][0])==type(True):
        page.drawString(250,500,head[2][1])
        page.close()
        exit()

    proc=head[-1][0]
    branch=head[-1][1]
    init=head[-1][2]
    sear=head[-1][3]

    ox=250
    oy=550
    sx=100
    sy=50
    sy=35

    page.drawString(ox,oy,"Drawing summary")
#
# draw initial one
#
    ll=0
    for i in proc[init]:
        _,w=getPropString(branch[i])
        ll+=w

    if ox+(ll+8)*(sx)>_magic_x:
        sx=(_magic_x-ox)/(ll+8)
#komai
    if ll>30:
        page.setFont(font_size=1.5)
    print("length",ll)

    page.drawString(ox+(2)*sx,oy-(1)*sy,"reactions")

#    ox=ox+sx*5
    ox=ox+sx*0
    yy=oy-sy*2

    first=True
    scale=20
    
    pos,deep=preDraw(init,proc,branch)
#
# the first line
#
    for i in pos:
        ix=pos[i][0]
        iy=pos[i][1]
        s,w=getPropString(branch[i])
        page.drawString(ox+(ix+2)*sx,oy-(iy+1)*sy,s)

    ddd={}
    for i in proc:
        first=True
        for j in proc[i]:
            x1=pos[j][0]
            y1=pos[j][1]
            if not first:
                draw=False
                if j0 in ddd:
                    if j in ddd[j0]:
                        ddd[j0].append(j)
                        draw=True
                else:
                    ddd[j0]=[j]
                    draw=True
#                if draw:
                if True:
                    page.drawArrow(arrow(ox+(x0+w-1)*sx,oy-(y0+1)*sy-4,ox+(x1+1.9)*sx,oy-(y1+1)*sy-4,scale,fc=True)) 
            s,w=getPropString(branch[j])
            first=False
            x0=x1
            y0=y1
            j0=j

    name={}
    for j in proc:
        if deep[j] in name:
            name[deep[j]]+=","+str(j)
        else:
            name[deep[j]]=str(j)

    ll=0
    for j in name:
        if ll<len(name[j]):
            ll=len(name[j])

    if ll>23:
        page.setFont(font_size=1)
    elif ll>10:
        page.setFont(font_size=1.5)
    else:
        page.setFont(font_size=2)

    bt=0
    for j in name:
        print(j,name[j])
        if bt<j:
            bt=j
        if len(name[j])>20:
            page.drawString(_magic_x+sx*2,oy-(j+1)*sy,name[j][:22])
            page.drawString(_magic_x+sx*2,oy-(j+1)*sy-sy/2,name[j][22:])
        else:
            page.drawString(_magic_x+sx*2,oy-(j+1)*sy,name[j])

    page.setFont(font_size=2)
    page.drawString(_magic_x+sx*3,oy-(1)*sy,"route #id")

    if oy-(bt+2)*sy-sy<0:
        py=oy-sy
    else:
        py=oy-(bt+2)*sy-sy

    page.drawString(ox,py,sear)
    py-=sy*0.5

    page.setFont(font_size=1)

    add="Routes start at similarity of "
    for i in sorted(hope):
        add+="{:.2f}".format(i[1])+"("+str(i[0])+"),"
        if len(add)>90:
            page.drawString(ox,py,add)
            py-=sy*0.01
            add=""

    if len(add)>0:
       page.drawString(ox,py-sy*0.5,add)

    page.stroke()
    page.setFont(font_size=1)

# komai2
def getList(s):
    ret=[]
    s=s.split('[')[1].split(']')[0].split(',')
    for i in s:
        ret.append(int(i))
    return ret

def getAllRoutes(all):
    ret=[]
    key='Route:'
    for f in all:
        if key in f:
#            print("--",f,"<--",f.split(key))
            num=int(f.split('\n')[0].split(key)[1])
            ret.append(num)

    return ret

def whichIsFirstMaterial(s,mode):
    #komai
    r=0;maxWeight=0.0
    result=s[-1]
    for i, smiles in enumerate(s):
        if ">" in smiles:
            break
        if mode==3:# number of atom
            mol=Chem.MolFromSmiles(smiles)
            mw=mol.GetNumAtoms()
        elif mode==2:# weigh of atom
#            mw=1-Levenshtein.ratio(smiles,result)
            mol=Chem.MolFromSmiles(smiles)
            mw=rdMolDescriptors.CalcExactMolWt(mol)
        else:# similarity between source and product
            mw=Levenshtein.jaro_winkler(result,smiles)
#        print(smiles,result,mw)
#        mol=Chem.MolFromSmiles(smiles)
#        mw=rdMolDescriptors.CalcExactMolWt(mol)
        if mw>maxWeight:
            maxWeight=mw
            r=i
    return r

def isIncluded(s,p):
    flag=False

    for i in p:
        if len(i)<len(s):
            continue
        flag=True
        for j in s:
            if not j in i:
                flag=False
                break
        if flag:
            return True
    return False

def getMatrix(rt,all):# remove redundant routes
    ret={}

    tg="smiles"
    ret[tg],ppt=getMatrixAgent(rt,all,tg)
    tg="products"
    ret[tg],_=getMatrixAgent(rt,all,tg)
    return ret,ppt

def isInRef(s,ref):
    if s in ref:
        return ref.index(s),False
    n=len(ref)
    ref.append(s)
    return n,True

def getMatrixAgent(rt,all,typ="smiles"):# remove redundant routes

    _smiles={}
    ref=[]
    sAll={}
    sTotalMatch={}
    sSubMatch={}
    sDrop=[]
    ret={}

    routes=rt.copy()

    for r in rt:
        _smiles[r]=getSmiles(r,all)

# small trick
    l=0
    pt=0
    for r in reversed(rt):
        ss=_smiles[r]
        if l<len(ss):
            l=len(ss)
            pt=r

    print("len",len(_smiles),l,pt)
    if l<1:
        return False

    for ss in reversed(_smiles[pt]):
        if typ=="smiles":
            ref.append(ss)
        else:
            ref.append(ss[-1])

    if typ!="smiles":
        ss=_smiles[pt][0][0]
        ref.append(ss)

#    print("Fist Ref",ref,sim)

    for r in reversed(rt):
#        print(r,end="->")
        m=[]
        if typ=="smiles":
            for ss in reversed(_smiles[r]):
                n,_=isInRef(ss,ref)
                m.insert(0,n)
        else:
            for ss in reversed(_smiles[r]):
                n,f=isInRef(ss[-1],ref)
                m.insert(0,n)

            s=_smiles[r][0][0]
            n,f=isInRef(s,ref)
            m.insert(0,n)
        sAll[r]=m
        
    print("type",typ)
    print("the trick",pt,sAll[pt])
    print(sAll)

    ppt=pt
    for r in rt:
        if not sAll[r] in sTotalMatch.values():
            sTotalMatch[r]=sAll[r]
            if sAll[r]==sAll[pt]:
                ppt=r

    ret["total"]=sTotalMatch

    sSubMatch=sTotalMatch.copy()

    for r in sTotalMatch:
        flag=False
        if not r in sSubMatch:
            continue
        for p in sTotalMatch:
            if p == r:
                continue
            if p not in sSubMatch:
                continue
            if len(sTotalMatch[r])>len(sTotalMatch[p]):
                continue
            flag=True
            for i in sTotalMatch[r]:
                if not flag:
                    break
                if not i in sTotalMatch[p]:
                    flag=False
            if flag:
                sSubMatch.pop(r)
                break

    ret["sub"]=sSubMatch

#        if type!="smiles":
#            print("hot now",r,m)

#    print("Total")
#    for r in sTotalMatch.keys():
#        print(r,sTotalMatch[r])

    if typ=="smiles":
        return ret,ppt

    print("this is submatch",sSubMatch)
    
    return ret,ppt

def getTheBottom(c,f):
    py=evalHeight(c,fc)
    return py-3

def getLeftTab():
    return 2

def setWds(c,fc):

    for n in range(len(c)):
        if n<len(fc):
            if len(c[n][2])<1:# starting point
                tg=-1
                for i in fc[n].keys():
                    if tg>0:
                        continue
                    if not i in c[n][-1]:
#                        print("keys",i,fc[n][i])
                        tg=fc[n][i]
                if tg<0:
                    fc[n]['s']=getLeftTab() # depend on
                else:
                    fc[n]['s']=tg # depend on
            else:
                p=c[n][2][0]#connecting node
                if p<0:# new tab
                    fc[n]['s']=getLeftTab() # depend on
                else:
                    q=c[p][1][0]# this is the target
                    fc[n]['s']=fc[p][q]
        else:
            fc.append({'s':2})

def isInIt(s,d):

    for i in s:
        if not i in d:
            return False
    return True

def registerBranch(p,rt):
    ret=[]
    for r in rt:
        if isInIt(p,rt[r]):
            ret.append(r)
    return ret

def checkIn(pp,rt):
    ret=[]
    if type(pp)==int:
        return checkIn2(pp,rt)

    for r in rt:
        if checkIn3(pp,rt[r]):
            ret.append(r)
#    print("checkIn",pp,ret)
    return ret

def checkIn3(p1,p2):

    ll=len(p1)
#    print("xxx",ll,p1,p2)

    for i in range(len(p2)-ll+1):
        flag=True
        for j in range(ll):
#            print(p2[i+j],"vs",p1[j])
            if p2[i+j]!=p1[j]:
                flag=False
        if flag:
            return flag
    return False

def checkIn2(p,rt):
    ret=[]
    for r in rt:
        if p in rt[r]:
            ret.append(r)
    return ret

def branchCheck(s,d):

    for i,j in enumerate(d):
#        print("<--",i,s,j,j==s)
        if j==s:
#            print("--->",i,j,s,index[i])
#            return index[i]
            return i
    return False

def deepArrange(routes):
    n=0
    p=[]
    branch=[]
    for r in routes.values():
        m=max(r)
        if n<m:
            n=m

    p=[n]
    n-=1
    while n>-1:
        p2=checkIn(p,routes)
        p1=checkIn(p+[n],routes)
        p3=checkIn([n],routes)
        if p1!=p2 or p1!=p3:
            branch.append(p)
            p=[n]
        else:
            p.append(n)
        n-=1
    branch.append(p)
    print(branch)

    ok={}
    for r in routes:
        pp=[]
        rr=[]
        for s in routes[r]:
            pp.append(s)
            key=branchCheck(pp,branch)
#            print(pp,branch,key)
#            print("->",pp,key,"in",routes[r])
            if not key is False:
                rr.append(key)
                pp=[]
        ok[r]=rr

    for i in ok:
        print("Route",i,end=":")
        for j in ok[i]:
            print(branch[j],end="->")
        print()
    return [ok,branch]

##############################################
#################### main ####################
##############################################

chemical="losartan"
proc_per_line=10 # how many reactions per line roughly ?
_fc=True
_tite=False
_type=3
_id=False
_product_only=False 
_include_subsets=False

#print(sys.argv)
ppt=False;chem=None;_routes=[]
input_dir="./"
com='-chem chemical_name, -ppt for power point -n [1,2,3] to specify routes -p proc_per_line -d input_path, -one_size : rendering at the same size, -product_only : categolize routes with respct to reaction product only, -show_subsets : show routes including subsets'
com+=" -summary 0/1/2"
ops={'-h':com}
skip=False
for n,op in enumerate(sys.argv):
    if skip:
        skip=False;continue
#    print(n,op)
    match op:
        case '-h':
            print(sys.argv[0],ops['-h']);exit()
        case '-ppt':
            ppt=True
        case '-heavy':
            _type=2
        case '-atom_number':
            _type=3
        case '-one_size':
            _fc=False
        case '-show_subsets':
            _include_subsets=True
        case '-n':
            _routes=getList(sys.argv[n+1])
            skip=True
        case '-tite':
            _tite=True
        case '-sub_id':
            _id=sys.argv[n+1]
            skip=True
        case '-d':
            input_dir=sys.argv[n+1]
            skip=True
        case '-product_only':
            _product_only=True
        case '-p':
            proc_per_line=int(sys.argv[n+1])
            skip=True
        case '-chem':
            chemical=sys.argv[n+1]
            skip=True

input_file=chemical+'.txt'
if _id:
    output_file=chemical+str(_id)
else:
    output_file=chemical

if ppt:
    output_file=output_file+'.pptx'
else:
    output_file=output_file+'.pdf'

if _web:
    output_file=input_dir+'/report/'+output_file

print(input_file)
print(output_file)
#print(_routes)
print(proc_per_line)

head=None

print(input_dir+"/"+input_file)

with open(input_dir+"/"+input_file,"r") as f:
    all=f.readlines()

allRoutes=getAllRoutes(all)
routeMatrix,ppt=getMatrix(allRoutes,all)

tg="smiles"
if _product_only:
    tg="products"
parent=routeMatrix[tg]

if _include_subsets:
    show='total'
    com="."
else:
    show='sub'
    com=" and subsets eliminated."

print("parent",parent['total'])
routes=list(parent[show].keys())

#easy="All output:"+str(len(allRoutes))+", Net output:"+str(len(parent['total']))+", Prime output:"+str(len(parent['sub']))

print(f"Categolized in respect to {tg}"+com)

for r in parent[show]:
    print(r,parent[show][r])

print("--->",parent[show][r])
if len(parent[show][r])<1:
    hint=[False,"No search results"]
else:
    hint=deepArrange(parent[show])
    hint.append(ppt)

span=proc_per_line*2
if span<10:
    span=10
scale=_magic_x/(span)

_py=500
py=_py
ox=40

if len(_routes)>0:
    routes=_routes

page=openReport(output_file,scale)


#print(head)

sim=similarityOnRoute(hint[2],hint[0],hint[1],all) 
easy=str(len(parent['total']))+" variations from "+str(len(sim))+" routes over "+str(len(allRoutes))+" queries"
hint.append(easy)
if parent!=False:
    head=[chemical,getSmiles(routes[0],all),hint]
else:
    head=[chemical,getSmiles(1,all),hint]
head_page(head,page,sim)

for route in routes:
    print(f"Route:{route}",end="..")
    oy=py-scale

    smiles=getSmiles(route,all)

    fc=makeSgv(smiles)

    if not _fc:
        fc=_fc

    connect,start=initConnect(smiles,fc=fc,mode=_type)
    Debug(connect,head='G',sp=5)
#
    align(connect,smiles,span-3,fc,_tite)
    Debug(connect,head='H',sp=5)
    hy=evalHeight(connect,fc)

    if fc:
        setWds(connect,fc)

    if oy+hy*scale*2<0:
        page.stroke()
        py=_py
        oy=py-scale
    page.setFont()

    tmp_goal=[]
    for i in range(len(connect)):
        x0=connect[i][4]*scale+ox
        y0=connect[i][5]*scale*2+oy

        if fc:
            wds=fc[i]['s']
        if len(connect[i][2])<1:
            fn=getName(i,connect[i][6]+1)# need
#            print("starting",i,fn,x0,y0)
            if not fc:
                page.drawImage(fn,x0,y0,scale) # for starting material
            else:
#                wds=fc[i][connect[i][6]]
                page.drawImage(fn,x0,y0,scale*wds) # for starting material
        n=connect[i][3]
#    if n<0:
#        continue
        if n<0:# for goal
            if not fc:
                xx= x0+(len(connect[i][7])+2)*scale
            else:
                if i<len(fc):
                    sp=len(connect[i][7])+2
                else:
                    sp=len(connect[i][7])+2

#                    xx=x0+(fc[-1]['cheat']+fc[i][connect[i][6]])*scale
                xx=x0+(wds+sp)*scale
            yy= y0
            if len(tmp_goal)<1 and n==-1:
                tmp_goal=[xx,yy]
        else:
            xx=connect[n][4]*scale+ox
            yy=connect[n][5]*scale*2+oy

        if connect[i][1][0]>-1:
#        if connect[i][3]>-1:
            fn=getName(i,connect[i][1][0]+1)
#            print("result",i,fn,xx,yy)
            if not fc:
                page.drawImage(fn,xx,yy,scale) # for resultant material
                xp=x0+scale*1.5
            else:
                page.drawImage(fn,xx,yy,scale*fc[i][connect[i][1][0]]) # for resultant material
                xp=x0+scale*(0.5+wds)

        for j in connect[i][-1]:
            fn=getName(i,j+1)
#            page.drawImage(fn,xp,y0+step*0.55,scale) # for catalytic material
            if not fc:
                page.drawImage(fn,xp,y0,scale,bt=True) # for catalytic material
                xp=xp+scale
            else:
#                page.drawImage(fn,xp,y0,scale*fc[i][j],bt=True) # for catalytic material
#                fx=fc[i][j] if fc[i][j]<2 else 1
                fx=1
                page.drawImage(fn,xp,y0,scale*fx,bt=True) # for catalytic material
                xp=xp+fc[i][j]*scale
#        page.drawArrow(arrow(x0,y0,xx,yy,scale)) 
        if n<0 and len(tmp_goal)>0:# for goal
            xx=tmp_goal[0]
            yy=tmp_goal[1]
## komai
#
        if not fc:
            page.drawArrow(arrow(x0,y0,xx,yy,scale)) 
        else:
            skip=wds
#            print("arrow:",n,wds,x0,xx,xx-x0,skip)
            page.drawArrow(arrow(x0,y0,xx,yy,scale,skip=skip,fc=fc))

        if fc:
            py=oy+getTheBottom(connect,fc)*scale
        else:
            if py>y0-scale:
                py=y0-scale

    page.drawString(ox,oy+scale*1.5,"Route"+str(route))

#for i in range(10):
#    page.drawLine([0,i,10,i])
page.close()
print("end")
