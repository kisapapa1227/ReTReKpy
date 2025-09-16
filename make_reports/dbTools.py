import sys,os,time
import math
import subprocess
from rdkit import Chem
from rdkit import DataStructs
from rdkit.Chem.Fingerprints import FingerprintMols
from rdkit.Chem import Draw, rdMolDescriptors
from rdkit.Chem.Draw import rdMolDraw2D
from IPython.display import SVG

from reportlab.graphics import renderPDF
from reportlab.pdfgen import canvas
from reportlab.pdfbase.pdfmetrics import stringWidth
from reportlab.lib.pagesizes import A4, portrait, landscape
from svglib.svglib import svg2rlg
from reportlab.lib.units import mm
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from PIL import Image
import shutil
import Levenshtein
import datetime

version="1.7 08052025"

_web=False
_web=True

if _web:
    dr='/var/www/html/public/images/tmp/' # working directory
else:
    dr='tmp/' # working directory

_magic_x=650
_magic_x=550

def topH():
    return 500

class openReport:
    def __init__(self,file,scale):
        self.left=210*mm*0.05
        self.top=297*mm*0.05
        self.height=210*mm*0.9
        self.width=297*mm*0.9
        self.scale=scale
        self.mg=20

        ext=file.split(".")[1]

        if ext=='txt':
            self.width*=1.5
            self.stack=[]

        if ext=='pdf' or ext=='txt':
            print("open pdf:",file)
            self.type=ext
            self.page=canvas.Canvas(file,pagesize=landscape(A4))
            self.font_name="Times-Bold"
            self.font_size=4*mm
            self._font_size=4*mm
            self.P2PRatio=1.0
        elif ext=='pptx':
            self.type='ppt'
            self.page=Presentation()
            self.page.slide_height=Inches(8.27)
            self.page.slide_width =Inches(11.69)
            blank_slide_layout = self.page.slide_layouts[6]
            self.slide = self.page.slides.add_slide(blank_slide_layout)
            self.file=file
            self.font_name="Times-Bold"
            self.font_size=12 # 14<-16 
            self._font_size=12 #
#            self.P2PRatio=self.width/297*mm
            self.P2PRatio=self.page.slide_width/(297*mm)

    def getSpan(self):
#       scale=_magic_x/(span)
        return _magic_x/self.scale
     
    def setFont(self,font_name=None,font_size=None):# size is relative value to default
        if self.type=='pdf':
            if font_name:
                self.font_name=font_name
            if font_size:
                self.font_size=int(self._font_size*font_size)
#                print("in setFont",self.font_size)
            self.page.setFont(self.font_name,self.font_size)
        else:
            if font_size:
                self.font_size=self._font_size*font_size
    def stroke(self):
        if self.type=='pdf':
            self.page.showPage()
            self.page.setFont(self.font_name,self.font_size)
        elif self.type=='txt':
            y0,y1=2000.0,-1000.0
#komai
            for s in self.stack:
                match s[0]:
                    case 'drawImage':
                        if y0>s[3]:
                            y0=s[3]
                        if y1<s[3]+s[5]*0.4:
                            y1=s[3]+s[5]+0.4
                    case 'line':
                        if y0>s[2]:
                            y0=s[2]
                        if y1<s[2]:
                            y1=s[2]
                        if y0>s[4]:
                            y0=s[4]
                        if y1<s[4]:
                            y1=s[4]
            for s in self.stack:
            #    print(s)
                match s[0]:
                    case 'drawImage':
                        print(s[0]+":",s[1],'{:1g}'.format(s[2]),'{:1g}'.format(y1-s[3]-s[5]*0.4),s[4],s[5],s[6],s[7],end=";")
                    case 'line':
                        print(s[0]+":",'{:1g}'.format(s[1]),'{:1g}'.format(y1-s[2]),'{:1g}'.format(s[3]),'{:1g}'.format(y1-s[4]),end=";")
            self.stack=[]
            #print()
        else:
            blank_slide_layout = self.page.slide_layouts[6]
            self.slide = self.page.slides.add_slide(blank_slide_layout)

#komai
    def getImageSize(self,fn):
        if _web:
            ret=subprocess.run(['/usr/bin/identify',fn],capture_output=True,text=True).stdout
        else:
            ret=subprocess.run(['identify',fn],capture_output=True,text=True).stdout
        xy=ret.split(" ")[2].split("x")

        return float(xy[0]),float(xy[1])

    def drawImage(self,fn,x,y,cat=False):
        if 'branch' in fn:
            return 0,0

        ppx=self.left+x*self.scale;ppy=self.top+y*self.scale
        if self.type=='pdf':
            drawing=svg2rlg(fn)
            drawing.renderScale=self.scale
            renderPDF.draw(drawing, canvas=self.page,x=ppx,y=ppy)
        elif self.type=='txt':
#            drawing=svg2rlg(fn)
#            print(fn)
#            drawing.renderScale=self.scale
#            self.stack.append(["drawImage",fn.split('/')[-1],ppx,ppy,drawing.width,drawing.height])
            width,height,smiles=self.sizeFromFile(fn)
            if cat:
                self.stack.append(["drawImage",fn.split('/')[-1],ppx,ppy,width,height,smiles,'cat'])
            else:
                self.stack.append(["drawImage",fn.split('/')[-1],ppx,ppy,width,height,smiles,'sub'])
            return 0,0
        else:
            key="2.0px"
            kkk="4.0px"

            with open(fn,"r") as f:
                all=f.readlines()

            gn=fn.split(".svg")[0]+".jpg"
            mag=360
            if _web:
                subprocess.run(['/usr/bin/convert','-density',f'{mag}',fn,gn])
            else:
                subprocess.run(['convert','-density',f'{mag}',fn,gn])

            ix,iy=self.getImageSize(fn)

            if ix>iy:
                sx=1.0;sy=iy/ix
            else:
                sy=ix/iy;sy=1.0
            ppx=self.left+x*self.scale;ppy=self.top+y*self.scale
            width=self.scale*self.P2PRatio*ix;height=self.scale*self.P2PRatio*iy
            x0,y0=ppx*self.P2PRatio,self.page.slide_height-(ppy)*self.P2PRatio-height
#            x0,y0=x*self.sc,self.page.slide_height-(y+iy)*self.sc
#            width=self.scale*self.P2PRatio*0.9*sx;height=self.scale*self.P2PRatio*0.9*sy
            self.slide.shapes.add_picture(
                gn,left=x0,top=y0,
                width=width, height=height
    )
#            printF(gn,x0,y0,width,height)
#            printF("(x,y):",x,y)
#        return ix,iy


    def sizeFromFile(self,fn):
        tg=fn.split('/')[-1]
        index="/".join(fn.split('/')[:-1])+"/svgFile.info"
        with open(index) as fp:
            for l in fp:
                if tg in l:
                    x=l.split('\n')[0].split(":")
                    return int(x[-3]),int(x[-2]),x[-1]
        exit()

    def drawLine(self,pt):
        x0=self.left+pt[0]*self.scale;y0=self.top+pt[1]*self.scale
        x1=self.left+pt[2]*self.scale;y1=self.top+pt[3]*self.scale
        if self.type=='pdf':
            self.page.line(x0,y0,x1,y1)
        elif self.type=='txt':
            self.stack.append(["line",x0,y0,x1,y1])
        else:
            x0,y0=x0*self.P2PRatio,self.page.slide_height-y0*self.P2PRatio
            x1,y1=x1*self.P2PRatio,self.page.slide_height-y1*self.P2PRatio
            putLine(x0,y0,x1,y1,self.slide.shapes)

    def drawArrow(self,pt):
        x0=self.left+pt[0]*self.scale;y0=self.top+pt[1]*self.scale
        x1=self.left+pt[2]*self.scale;y1=self.top+pt[3]*self.scale
        x2=self.left+pt[4]*self.scale;y2=self.top+pt[5]*self.scale
        x3=self.left+pt[6]*self.scale;y3=self.top+pt[7]*self.scale

        if self.type=='pdf':
            self.page.line(x0,y0,x1,y1)
            self.page.line(x1,y1,x2,y2)
            self.page.line(x1,y1,x3,y3)
        elif self.type=='txt':
            self.stack.append(["line",x0,y0,x1,y1])
            self.stack.append(["line",x1,y1,x2,y2])
            self.stack.append(["line",x1,y0,x3,y3])
        else:
            x0,y0=x0*self.P2PRatio,self.page.slide_height-y0*self.P2PRatio
            x1,y1=x1*self.P2PRatio,self.page.slide_height-y1*self.P2PRatio
            x2,y2=x2*self.P2PRatio,self.page.slide_height-y2*self.P2PRatio
            x3,y3=x3*self.P2PRatio,self.page.slide_height-y3*self.P2PRatio
            group=[]
            group+=[putLine(x0,y0,x1,y1,self.slide.shapes)]
            group+=[putLine(x1,y1,x2,y2,self.slide.shapes)]
            group+=[putLine(x1,y1,x3,y3,self.slide.shapes)]
            self.slide.shapes.add_group_shape(group)

    def getProperLength(self,mag,string):
        if self.type=='pdf' or self.type=='txt':
            tx=string
            while True:
                t=(self.getSpan()-1)*self.scale-stringWidth(tx,self.font_name,self.font_size)
                if t>0:
                    return len(tx)
                tx=tx[0:len(tx)-3]
        else:
            return 50

    def drawStringR(self,x,y,string,center=False,adj=True):
        xx=self.left+x*self.scale
        if center:
            xx-=stringWidth(string,self.font_name,self.font_size)/2
        yy=self.top+y*self.scale
        if self.type=='pdf':
            self.page.drawString(xx,yy,string)
        elif self.type=='txt':
            self.stack.push(["string:",xx,yy,string])
        else:
            xx=xx*self.P2PRatio
            yy=self.page.slide_height-yy*self.P2PRatio-self.font_size
            if adj:
                pw=Cm(0.1)*self.font_size*(len(string)+2)*0.20
                py=Cm(0.1)*self.font_size*0.5
                shape=self.slide.shapes.add_textbox(xx,yy-Pt(self.font_size),pw,py)
            else:
                pw=self.page.slide_width*0.9
                sw=self.font_size*(len(string)+2)*0.18/pw+1
                shape=self.slide.shapes.add_textbox(xx,yy-Pt(self.font_size),pw,py)

            tf=shape.text_frame
            tf.word_wrap=True
#            tf.auto_size=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
            p=tf.paragraphs[0]
            run=p.add_run()
            run.text=string
            run.font.name=self.font_name
            run.font.size=Pt(self.font_size)
            sw=self.font_size*(len(string)+2)*0.18/pw+1
            shape=self.slide.shapes.add_textbox(xx,yy-Pt(self.font_size),pw,py)

    def close(self):
        if self.type=='pdf':
            self.page.save()
        else:
            self.page.save(self.file)
    def isPdf(self):
        if self.type=='pdf':
            return True
        return False

def fm(x):
    return "{:.1f}".format(x)

def putLine(x0,y0,x1,y1,shapes):
    if x1>x0:
        if y1>y0:
           shape=shapes.add_shape(MSO_SHAPE.LINE_INVERSE,x0,y0,x1-x0,y1-y0)
           shape.element.getchildren()[1].getchildren()[0].set('flipV','1')
        else:
           shape=shapes.add_shape(MSO_SHAPE.LINE_INVERSE,x0,y1,x1-x0,-y1+y0)
    else:
        if y1>y0:
           shape=shapes.add_shape(MSO_SHAPE.LINE_INVERSE,x1,y0,-x1+x0,y1-y0)
        else:
           shape=shapes.add_shape(MSO_SHAPE.LINE_INVERSE,x1,y1,-x1+x0,-y1+y0)
           shape.element.getchildren()[1].getchildren()[0].set('flipV','1')
    shape.fill.solid()
    shape.fill.fore_color.rgb=RGBColor(0,255,0)
    shape.line.color.rgb=RGBColor(0,0,0)
    return shape

def cropSvg(fn):
    if not os.path.isfile(fn):
        return [0,0]

    fno=fn+".tmp";fo=open(fno,"w")

    with open(fn,"r") as fp:
        for l in fp:
            fo.write(l)
            if 'viewBox' in l:
                w=l.split("width='")[1].split('px')[0]
                h=l.split("height='")[1].split('px')[0]

    x=[float(w),0]
    y=[float(h),0]
    fo.close()
    inPath=False
    with open(fn,"r") as fp:
        for l in fp:
            if "d='" in l:
                inPath=True
            if inPath:
                if "d=" in l:
                    if '/>' in l:
                        inPath=False
                    l=l.split("d=")[1].split("'")[1].split("\n")[0]
                f,v=extF(l)
                if f:
                    px=[]
                    for p in v:
                        try:
                            px.append(float(p))
                        except:
                            px.append(0.0)
                        if len(px)==2:
                            repv(x,y,px)
                            px=[]
                else:
                    inPath=False
    mg=10

    x0=int(x[0]-mg);x1=int(x[1]+mg)
    y0=int(y[0]-mg);y1=int(y[1]+mg)

    sx=int(x1-x0)
    sy=int(y1-y0)

    ff=open(fn,"w")
    with open(fno,"r") as fp:
        for l in fp:
            if "<rect " in l:
                ff.write(l)
            elif "stroke-width:" in  l:
                ff.write(l.replace("stroke-width:2.0","stroke-width:1.0"))
            elif "viewBox" in  l:
                ff.write(f"width='{sx}px' height='{sy}px' viewBox='{x0} {y0} {sx} {sy}'>\n")
            else:
                ff.write(l)
    ff.close()
    os.remove(fno);

    return [sx,sy]

def chkProc(dr,name):
    with open(dr+"/svgFile.info") as f:
        lines = f.readlines()
    fc=-1;rs=False
    for line in lines:
        if name in line:
            r=line.split("\n")[0].split(":");
            fc=float(r[1]);rs=[int(r[2]),int(r[3])]
    return fc,rs

def smile2svg(smiles,name,dr=dr,size=200,newProc=False):
#    if os.path.exists(dr):
#        shutil.rmtree(dr)

    if not os.path.exists(dr):
        os.makedirs(dr)
        with open(dr+"/svgFile.info",'w') as f:
            f.write("this file is made by readDb.py\n");

    if newProc:
        fc,r=chkProc(dr,name)
        if fc>0:
            return fc,r
    
    mol=Chem.MolFromSmiles(smiles)
    mw=mol.GetNumAtoms()
    fc=int(math.sqrt(mw)/1.2)

    if fc<1:
        fc=1

    view = rdMolDraw2D.MolDraw2DSVG(size*fc,size*fc)

    option=view.drawOptions()
    option.circleAtoms=False
    option.continuousHightlight=False

    tm = rdMolDraw2D.PrepareMolForDrawing(mol)
    view.DrawMolecule(tm)
    view.FinishDrawing()

    svg=view.GetDrawingText()

    with open(dr+"/"+name,'w') as f:
        f.write(svg)

    r=cropSvg(dr+"/"+name)

    with open(dr+"/svgFile.info",'a') as f:
        f.write(name+":"+str(fc)+":"+str(r[0])+":"+str(r[1])+":"+smiles+"\n")
    return fc,r

#    os.remove(fo)

def Debug(c,head=None,sp=None):
    return
    for n,i in enumerate(c):
        if head:
            print(head,end=":")
        if sp:
            print(n,i[sp],end="<--")
        print(n,i)
#    exit()

def printF(head,*args):
    print(head,end='->')
    for arg in args:
        print(round(arg,2),end=" ")
    print()

def makeSvg(smiles,dst,type=0):
    size=100;fsize={};fcs=[];retSmiles=[]
    for n,proc in enumerate(smiles):
        fc={}
        for m,comp in enumerate(proc):
            if ">" in comp:
                continue
            fn="e"+str(n+1)+"x"+str(m+1)+".svg"
            fc[m],fsize[fn]=smile2svg(comp,fn,dr=dst,newProc=True,size=size)
        fcs.append(fc)
    with open(dst+"/svgFile.info") as f:
        lines = f.readlines()
    if not 'size' in lines[-1]:
        dd=datetime.datetime.now().strftime('%Y%m%d%H%M%S')
        ret=subprocess.run(['du',dst],capture_output=True,text=True).stdout
        with open(dst+"/svgFile.info",'a') as f:
            f.write('date:'+dd+"\n")
            f.write('size:'+ret.split('\t')[0]+"\n")
        
    return fcs,fsize

def arrowR(x0,y0,x1,y1):

    l=math.sqrt((x1-x0)*(x1-x0)+(y1-y0)*(y1-y0))
    sx=50
    pi=3.1415/12.0 # 30 deg

    x2=x0;y2=y0
    x3=x1;y3=y1

    a=math.atan2(y1-y0,x1-x0)
    x4=x3-sx*math.cos(a+pi);y4=y3-sx*math.sin(a+pi)
    x5=x3-sx*math.cos(a-pi);y5=y3-sx*math.sin(a-pi)

    return [x2,y2,x3,y3,x4,y4,x5,y5]

def arrow(x0,y0,x1,y1,scale,skip=1.0,fc=False,ugly_patch=True):

    if not fc:
        return prevArrow(x0,y0,x1,y1,scale,ugly_patch=ugly_patch)
#
# ugly patch
#
    if ugly_patch and x0>x1:
        x1=x0+scale*3
        y1=y0
# end
# anyway
    l=math.sqrt((x1-x0)*(x1-x0)+(y1-y0)*(y1-y0))
    rt=(l-scale*skip)/l
    sx=10
    pi=3.1415/12.0 # 30 deg

#from: shift full
    x2=x0+skip*scale
    y2=y0+scale/2
#to: shift full 
    x3=x1
    y3=y1+scale/2

    a=math.atan2(y3-y2,x3-x2)

    xx=x3-sx*math.cos(a+pi);yy=y3-sx*math.sin(a+pi)
    ret=[x2,y2,x3,y3,xx,yy]

    xx=x3-sx*math.cos(a-pi);yy=y3-sx*math.sin(a-pi)

    return ret+[xx,yy]

def prevArrow(x0,y0,x1,y1,scale,skip=0.0,ugly_patch=True):
#
# ugly patch
#
    if ugly_patch and x0>x1:
        x1=x0+scale*3
        y1=y0
# end

    l=math.sqrt((x1-x0)*(x1-x0)+(y1-y0)*(y1-y0))
    rt=(l-scale*skip)/l
    sx=10
    pi=3.1415/12.0 # 30 deg

    x2=x1+(x0-x1)*(rt+0.10);y2=y1+(y0-y1)*(rt-0.10)
    x3=x0+(x1-x0)*(1.0+0.10);y3=y0+(y1-y0)*(1.0-0.10)

    a=math.atan2(y1-y0,x1-x0)
    x4=x3-sx*math.cos(a+pi);y4=y3-sx*math.sin(a+pi)
    x5=x3-sx*math.cos(a-pi);y5=y3-sx*math.sin(a-pi)

    return [x2,y2,x3,y3,x4,y4,x5,y5]

def getName(a,b,dr=dr):
    fn=dr+"e"+str(a+1)+"x"+str(b)+".svg"#<-- need modification
    return(fn)

def getGoal(smiles,connect):
    goal=[]
    l=len(smiles)
    for i in range(l):
        if connect[i][3]==-1:
            goal.append(i)
    return goal

def _ccn(i):
    if i<10:
        return 1
    if i<100:
        return 2
    return 3

def getPropString(page,m):
    if len(m)<2:
        s=str(m[0])
        w=stringWidth(s,page.font_name,page.font_size)*page.scale
        return s,w
    s=str(m[0])+"-"+str(m[-1])
    w=stringWidth(s,page.font_name,page.font_size)*page.scale
    return s,w

def easyTrans(page,x,y,half=False,lf=1.2):
    nx,ny=x/page.scale,y/page.scale+(page.height-page.font_size)/page.scale
    if half:
        ny+=page.font_size/page.scale/2.0
    return nx,ny,y-page.font_size*lf

def toStr(head):
    if len(head)<2:
        return str(head[0])
    return str(head[0])+"-"+str(head[-1])

def head_page(head,page,arrange=True):
    ox=40;oy=480;mag=3

    nx,ny,_=easyTrans(page,0,0)
    ll=len(head[0])
    page.setFont(font_size=mag)
    if ll>7:
        page.setFont(font_size=mag*7.0/float(ll))
    page.drawStringR(nx,ny,head[0])
    if ll>7:
        page.setFont(font_size=mag)

    scale=page.scale
    tmp_scale=page.scale
#
# head
#
    smiles=head[1][-1][-1]
    fn="head.svg"
    x0=ox
    smile2svg(smiles,fn,size=100)
    scale=page.scale
    page.scale=0.7
    print("dr<--------------",dr)
    page.drawImage(dr+"/"+fn,x0,page.top)

    page.scale=scale

    pt=0;l=len(smiles)
#    py=py-mag*page.scale
    mg=2
    page.setFont(font_size=mg)
    wd=page.getProperLength(mg,smiles)
#    wd=7
    tl=smiles
    py=-page.height
    while pt<l:
        end=min(pt+wd,l)
        tl=smiles[pt:end]
        nx,ny,_=easyTrans(page,page.width/2,py)
        page.drawStringR(nx,ny,tl,center=True)
        pt+=wd
        py+=page.font_size

    if type(head[2])==type(True):# no route
        page.setFont(font_size=6)
        nx,ny,_=easyTrans(page,page.width/2,-page.height/3)
        page.drawStringR(nx,ny,"Route Not Found.",center=True)
        return

    if type(head[2][0])==type(True):
        nx,ny,_=easyTrans(page,page.width/2,-page.height/2)
        page.drawStringR(nx,ny,head[2][1])
        page.close()
        with open(log_name,"a") as fp:
            fp.write(f"mission ok\n")
        exit()

    oy=550;sy=50;sy=35

#    print("font",page.font_size,page.scale)
    py=0
# here
    nx,ny,py=easyTrans(page,page.width/2,py)
    page.drawStringR(nx,ny,"Reaction route summary",center=True)

    page.scale=tmp_scale

    py-=page.font_size
    nx,ny,py=easyTrans(page,page.width*0.2,py)
# here : * variations from * routes over * queries
    page.drawStringR(nx,ny,head[-1][-1])

    py-=page.font_size
#
# draw initial one # ll=0
    page.scale=tmp_scale
    bt=0
    sx=14
    if len(head[2][0])>10:
        p2_font_size=mg*10.0/float(len(head[2][0]))
        page.setFont(font_size=p2_font_size)
        page.scale*=10.0/float(len(head[2][0]))
        p2_scale=page.scale

    keepPy=py
    nx,ny,py=easyTrans(page,page.width*0.8,py)
    page.drawStringR(nx,ny,"route #id")
    
    p=1
    for n,route_id in enumerate(head[2][0]):
#        print("the ",route_id)
        nx,ny,py=easyTrans(page,page.width*0.8,py)
        if arrange:
            ids,p=setIds(p,route_id[0])
        else:
            ids=setOrigIds(route_id[0])
        page.drawStringR(nx,ny,ids)
#        page.drawStringR(nx,ny,toStr(route_id[0]))

    py=keepPy
    mes="reactions"
    nx,ny,py=easyTrans(page,page.width*0.2,py)
    page.drawStringR(nx,ny,mes)

#    p2=head[2][2]
#    for ps in head[2][0]:
#        for p in ps[1]:
#            print(p2[p],end="->")
#        print("ok")

    pos=preDraw(page,head[2])
    keep_py=py

    for n in pos.keys():
        pp=py+page.font_size*1.2*pos[n][2]
        nx=page.width*0.7/page.scale-page.font_size*pos[n][1]*2
        ny=pp/page.scale+(page.height-page.font_size)/page.scale

        page.drawStringR(nx,ny,pos[n][0])

    y=0
    for ps in head[2][0]:
        p=ps[1][0]
        if pos[p][2]!=y:
            pp=py+page.font_size*1.2*y
            nx=page.width*0.7/page.scale-page.font_size*pos[p][1]*2
            ny=pp/page.scale+(page.height-page.font_size)/page.scale
            page.drawStringR(nx,ny,pos[p][0])

        y=y-1

    y=0
    for ps in head[2][0]:
        flag=False
        p=ps[1][0]

        pp=py+page.font_size*1.2*y
        nx1=page.width*0.7/page.scale-page.font_size*pos[p][1]*2
        ny1=pp/page.scale+(page.height-page.font_size)/page.scale
        ll=len(str(pos[p][0]))+1

        for i in range(len(ps[1])-1):
            p=ps[1][i+1]
            pp=py+page.font_size*1.2*pos[p][2]
            nx=page.width*0.7/page.scale-page.font_size*pos[p][1]*2
            ny=pp/page.scale+(page.height-page.font_size)/page.scale
            #komai
            page.drawArrow(
                    arrowR(nx1+ll*page.font_size*2,
                        ny1+page.font_size,
                nx,ny+page.font_size))
            nx1=nx;ny1=ny
            ll=len(str(pos[p][0]))+1
            flag=True
        y=y-1

    page.stroke()
    page.setFont(font_size=1)

def preDraw(page,proc):

    p2=proc[2]

#    for ps in proc[0]:
#        for p in ps[1]:
#            print(p2[p],end="->")
#        print("ok")

    d=[-1]*len(proc[2]) 
    pos={}

    y=0;x=0
    for p in reversed(proc[0][0][1]):
        if d[p]<0:
            s=toStr(p2[p])
            d[p]=0;x+=len(s)+4;pos[p]=[s,x,y]
#            print(s,x,y,p)

    x=0;y=-1
    for i in range(len(proc[0])-1):
        flag=True
        for p in reversed(proc[0][i+1][1]):
            if d[p]<0:
                s=toStr(p2[p])
                d[p]=0;px=px+len(s)+2;pos[p]=[s,px,y]
            else:
                s=toStr(p2[p])
                px=pos[p][1]
        y-=1

    return pos

def setOrigIds(route_id):
    ret=str(route_id[0])
    for i in range(1,len(route_id)):
        ret=ret+","+str(route_id[i])

    return ret

def setIds(p,route_id):
    ll=len(route_id);r=str(p)

    if ll>1:
        r=r+"-"+str(p+ll-1)

    return r,p+ll

def getList(s):
    ret=[]
    s=s.split('[')[1].split(']')[0].split(',')
    for i in s:
        ret.append(int(i))
    return ret

def whichIsFirstMaterial(s,mode):
    r=0;maxWeight=0.0
    result=s[-1]
    for i, smiles in enumerate(s):
        if ">" in smiles:
            break
        if mode==3:# number of atom
            mol=Chem.MolFromSmiles(smiles)
            mw=mol.GetNumAtoms()
        elif mode==2:# weigh of atom
#            mw=1-Levenshtein.ratio(smiles,result)
            mol=Chem.MolFromSmiles(smiles)
            mw=rdMolDescriptors.CalcExactMolWt(mol)
        else:# similarity between source and product
            mw=Levenshtein.jaro_winkler(result,smiles)
        if mw>maxWeight:
            maxWeight=mw
            r=i
    return r

def divList(src):
    ret=[]
    r=[src[0]]
    for i in range(len(src)-1):
        if src[i]-1==src[i+1]:
            r.append(src[i+1])
            continue
        else:
            ret.append(r)
            r=[src[i+1]]
    ret.append(r)
    return ret

def listMatch(s1,ss):

    for s2 in ss:
        if len(s1)==len(s2) and s1[0]==s2[0]:
            return True
    return False

def cutFlagment(src,key,i):
    rs=[];r=[]
    nk=key.index(i);ns=src.index(i) 
    lk=len(key);ls=len(src)

#    print("in cutFlagment",ns,ls-ns,lk-nk)
#    print(src,key)
    if ns>0:
        for i in range(ns):
            r.append(src[i])
        rs.append(r);r=[]

#    print(lk+ns)
    if (ls-ns)>(lk-nk):
        for i in range(lk+ns,ls):
            r.append(src[i])
        rs.append(r);r=[]

    return rs 

def cutC(s,cs):
    r=[];flag=False
    for c in cs:
        if len(s)>len(c):
            for i in c:
                if i in s:
                    r=cutFlagment(s,c,i);flag=True
                    return r,flag
    return r,flag

def getFullFlagment(cv):

    needIt=True
#    print("getFullFlag",cv)
    while needIt:
        needIt=False
        for n,c in enumerate(cv):
            rs,flag=cutC(c,cv)
            if flag:
                cv.pop(n);
                for r in rs:
                    if not listMatch(r,cv):
                        cv.append(r)
                needIt=True
    return cv

def findInFlagment(s,flags):
    for n,f in enumerate(flags):
        if s in f:
            return n
    return -1

def recomposeRoute(src,flags):
    ret=[]
    n=-1;
    for s in src:
        t=findInFlagment(s,flags)
        if t!=n:
            ret.append(t);n=t
    return ret

def getReactionSummary(route,idx):
    ok=[];cv=[]

    for r in route:
        cv.append(divList(r[0]))

    for c1 in cv:
        for c in c1:
            if not c in ok:
                ok.append(c)

    flagments=getFullFlagment(ok)
    ok=[];rr=[]
    for r in route:
        rc=recomposeRoute(r[0],flagments)
        p1=[]
        for x in r[1]:
#            p1.append(idx.index(x)+1)
            p1.append(x)
        ok.append([p1,rc]);rr.append(rc)

#    for o in ok:
#        print(o[0],end=";")
#        for r in o[1]:
#            print(flagments[r],end=";")
#        print("<--")

# rearrange by starting matrials and number of reaction nsteps
    g=makeSortIndex1(rr)
    ok,idx,route_map=makeSortIndex3(ok,idx,g)
#
    return [ok,idx,flagments,route_map]

def makeSortIndex3(ok,idx,g):
    ds=[-1]*len(ok)
    rrr=[];iii=[];ret=[]

    for l,d in enumerate(ds):
        n=g[l]
        if ds[n]>0:
            continue
        rrr.append(ok[n]);iii.append(n+1)
        start=ok[n][1][0];ds[n]=1
        for m in range(len(ds)):
            if ds[m]>0:
                continue
            if ok[m][1][0]==start:
                ds[m]=1;rrr.append(ok[m]);iii.append(m+1)
            m=m+1
    for r in rrr:
        ret=ret+r[0]

    return rrr,iii,ret

def makeSortIndex2(ok,idx,g):
    ret=[];iii=[];m=1
    for n in g:
        r=[];ii=[]
        for o in ok[n][0]:
            r.append(m);m+=1
            r.append(m);m+=1
#            ii.append(idx[o-1])
            ii.append(o)
        ret.append([r,ok[n][1]])
    return ret,ii

def makeSortIndex1(rc):

    dst=[-1]*len(rc);ret=[]
    for i in range(len(rc)):
        l=0;g=0
        for n,d in enumerate(dst):
            if d>-1:
                continue
            if len(rc[n])>l:
                l=len(rc[n]);g=n
        ret.append(g);dst[g]=0
        if l==0:
            break

    return ret

def strToParent(ss):
    ret={}
    n={1:'total',2:'sub'}

    for i in range(1,3):
        r={}
        for s in ss[i].split(';')[0:-1]:
            key=int(s.split(':')[0])
            values=s.split(':')[1]
            p=[]
            for q in values.split(',')[0:-1]:
                p.append(int(q))
            r[key]=p
        ret[n[i]]=r
    return ret

def strToSmiles(ss):
    ret=[]
    for s in ss.split('##')[0:-1]:
        item=[]
        for p in s.split(';')[0:-1]:
            item.append(p)
        ret.append(item)
    return ret

def strToList(ii):
        w=[]
        for i in ii.split(',')[0:-1]:
            w.append(int(i))
        return w

def strToConnect(ss):
    ret=[]
    for s in ss.split('##')[0:-1]:
        item=[]
        ii=s.split(';')
        item.append(ii[0])
        item.append(strToList(ii[1]))
        item.append(strToList(ii[2]))
        for n in range(3,7):
            item.append(int(ii[n]))
#        item.append('A')
        item.append(strToList(ii[7]))
        ret.append(item)
    return ret

def dump(fall,pid,odir):
    with open(odir+str(pid)+'db.txt','w') as fp:
        for xx in fall:
            fp.write('route '+str(xx[0])+'\n')
            fp.write(xx[2]+'\n')
            fp.write(xx[3]+'\n')

def repv(x,y,p):
    if x[0]>p[0]:
        x[0]=p[0]
    if x[1]<p[0]:
        x[1]=p[0]
    if y[0]>p[1]:
        y[0]=p[1]
    if y[1]<p[1]:
        y[1]=p[1]

def frac(ss,k):
    rr=[]
    for s in ss:
        if k in s:
            tt=s.split(k)
            for t in tt:
                if not t=="":
                    rr.append(t)
        else:
            rr.append(s)
    return rr

def xsplit(v,s):
    vv=[]
    for rx in v:
        for v in rx.split(s):
            if v!='':
                vv.append(v)
    return vv

def extF(l):
    if "class='atom" in l:
        return False,[]

    v=frac([l],"M")
    v=frac(v,"Q")
    v=frac(v,"L")

    vv=xsplit(v,',')
    vvv=xsplit(vv,' ')
    if "/>" in vvv[-1]:
        return False, False
    return True,vvv

def getSubObj(page,n,c,fsize,xx,onset=False):

    ret={}
    if onset:
        fn=getName(n,c[n][6]+1,dr='')
    else:
        fn=getName(n,c[n][1][0]+1,dr='')

    if fn in fsize:
        r=fsize[fn]
    else:
        r=[page.mg,page.mg]

    xx-=page.mg+r[0]
    ret['type']=1;ret['pos']=[[xx,-r[1]/2.0]];ret['svg']=[fn];ret['size']=r
    return ret,xx

def calcBranch(obj,tag):
    base=0
    x=obj[tag]['pos'][0]
    y=obj[tag]['pos'][1]+obj[tag]['size'][1]/2.0

    for o in obj:
        tx=o['pos'][0]
        ty=o['pos'][1]
        if tx<x:
            if base>ty:
                base=ty
    return x,base

def getCatObj(page,n,c,fsize,xx):
    l=0
    fns=[]
    height=0
    ret={}

    for j in c[n][-1]:# catalytic materials
        fn=getName(n,j+1,dr='')
        fns.append(fn)
        r=fsize[fn]
# will be check nicely
        h=(page.mg+r[1])
        if height<h:
            height=h
        l+=(page.mg+r[0])
    l+=(page.mg)
    if len(c[n][-1])<1:
        l+=page.mg*5
#Need..
    xx-=page.mg
    ret['type']=2;ret['pos']=[[xx-l,0],[xx,0]];ret['svg']=fns;ret['size']=[l,height]
    return ret,xx-l

def setItem(a,b,c,d):
    ret={}
    ret['type']=a;ret['pos']=b;ret['svg']=c;ret['size']=d
    return ret

def shiftForBranch(objs):

    for i in range(1,len(objs)):
        x0=objs[i][-1]['pos'][0][0]
        x1=objs[i][ 0]['pos'][0][0]+objs[i][ 0]['size'][0]
        ty=-topH();by=topH()
        for item in objs[i]:
            if item['type']==2:
                yy=item['pos'][0][1]+item['size'][1]
                if ty<yy:
                    ty=yy
        for item in objs[0]:
            if item['type']==1:
                p0=item['pos'][0][0]
                p1=item['pos'][0][0]+item['size'][0]
                if p0<x1 and x0<p1:
                    yy=item['pos'][0][1]-item['size'][1]/2
                    if by>yy:
                        by=yy
        shift=ty-by+100
        fool=0
        for item in objs[i]:
            if fool==0:
                item['pos'][0][1]-=shift
                fool+=1
                continue
            item['pos'][0][1]-=shift
            if item['type']==2:
                item['pos'][1][1]-=shift

def lineBraker(page,objs):
    for j in range(len(objs[0])-1,-1,-1):
        item=objs[0][j]
        xx=item['pos'][0][0]+item['size'][0]
        if xx>page.width/page.scale*0.9:
            cutIt(page,objs,item['pos'][0][0])

def cutIt(page,objs,ok):
    theWidth=page.width/page.scale*0.9
    by=False;ty=False
    for obj in objs:
        for item in obj:
            xx=item['pos'][0][0]+item['size'][0]
            if xx<theWidth:
                if item['type']==1:
                    yy=item['pos'][0][1]-item['size'][1]/2.0
                else:
                    yy=item['pos'][0][1]

                if by==False:
                    by=yy
                else:
                    if by>yy:
                        by=yy
#            else:
            elif xx<theWidth*2:
                if item['type']==2:
                    yy=item['pos'][0][1]+item['size'][1]
                else:
#
# corrected 08192025
#                    yy=item['pos'][0][1]+item['size'][1]/2.0 orig
#                    yy=item['pos'][0][1]+item['size'][1]/2.0
                    yy=item['pos'][0][1]+item['size'][1]/2.0
                if ty==False:
                    ty=yy
                else:
                    if ty<yy:
                        ty=yy
    shift=ty-by+page.mg
    for obj in objs:
        for item in obj:
            xx=item['pos'][0][0]+item['size'][0]
            if xx>theWidth:
                item['pos'][0][0]-=(ok-page.width*0.1)
                item['pos'][0][1]-=shift
                if item['type']==2:
                    item['pos'][1][0]-=(ok-page.width*0.1)
                    item['pos'][1][1]-=shift

def easyGetBottom(objs):
    hy=topH()
    top=-topH()
    for obj in objs:
        for item in obj:
            if item['type']==1:
                yy=item['pos'][0][1]-item['size'][1]/2.0
                if hy>yy:
                    hy=yy
            else:
                yy=item['pos'][0][1]+item['size'][1]
                if top<yy:
                    top=yy
    return top-hy,top,hy

def stackGoal(connect,n,line,tag,i=1):

    ret=[connect[n][2][i],line,tag]
    if len(connect[n][7])>0:
        ret[2]+=1

    return ret

def getNext(c,tmp_goal):
    if len(c)>0:
        return c[0]
    return -1

def makeObject(page,connect,fsize):

    line=0
    objs=[];hint=[];tmp_goal=[];xx=0
#
# find goal
#
    n=0
#    connect[13][3]=6
#    print(connect,len(connect))
    for i in range(len(connect)):
        if connect[i][3]<0:
            goal=n=i
            break
#
# temporal base y=0
#
    tag=0;obj=[]
#    fn=getName(n,connect[n][6]+1,dr='') # goal
    hint.append([]) # to which the line connect (line,tag)

    try:
        l=len(connect[n][2])
    except:
        l=0

    if l>1:
        tmp_goal.append(stackGoal(connect,n,line,tag))

    while n>-1:
        item,xx=getSubObj(page,n,connect,fsize,xx);obj.append(item);tag+=1
        l=len(connect[n][2])
        if l>1:
            for i in range(1,l):
                tmp_goal.append(stackGoal(connect,n,line,tag+i,i=i))

        item,xx=getCatObj(page,n,connect,fsize,xx);obj.append(item);tag+=1

        if l<1:
            item,xx=getSubObj(page,n,connect,fsize,xx,onset=True);obj.append(item);tag+=1
            objs.append(obj);line+=1
            break
        n=getNext(connect[n][2],tmp_goal)

    while len(tmp_goal)>0:
        r=tmp_goal.pop()
        n=r[0];obj=[];tag=0
#        hint.append([r[1],r[2]]) # to which the line connect (line,tag)
#        item=objs[r[1]][r[2]] # original 20250904
        item=objs[r[1]][r[2]-1] # connecting position
#        x0=item['pos'][0][0]-page.mg
        x0=item['pos'][0][0]-page.mg
        y0=item['pos'][0][1]+item['size'][1]/2
#        item=objs[r[1]][r[2]+2]
        ln=item['size'][0]
#        xx=x0-ln
        xx=x0-ln
        item=setItem(2,[[xx,0],[x0,y0]],[],[ln,0]);
        obj.append(item);tag+=1

        n=getNext(connect[n][2],tmp_goal)
        if n>-1:
            item,xx=getSubObj(page,n,connect,fsize,xx);obj.append(item);tag+=1
            item,xx=getCatObj(page,n,connect,fsize,xx);obj.append(item);tag+=1

        if len(connect[n][2])<1:
            item,xx=getSubObj(page,n,connect,fsize,xx,onset=True);obj.append(item);tag+=1
            objs.append(obj)
            n=-1
        else:
            n=getNext(connect[n][2],tmp_goal)

        while n>-1:
# catalysis object
            l=len(connect[n][2])
            if l>1:
                for i in range(1,l):
                    tmp_goal.append(stackGoal(connect,n,line,tag+i,i=i))

            item,xx=getCatObj(page,n,connect,fsize,xx);obj.append(item);tag+=1
            item,xx=getSubObj(page,n,connect,fsize,xx);obj.append(item);tag+=1
# up to now
            if l<1:
                objs.append(obj);line+=1
                break
                n=getNext(connect[n][2],tmp_goal)

    xx=0
    for obj in objs:
        for item in obj:
            if xx>item['pos'][0][0]:
                xx=item['pos'][0][0]

    for obj in objs:
        for item in obj:
            for i in range(len(item['pos'])):
                item['pos'][i][0]-=xx

    if len(objs)>0:
        shiftForBranch(objs)
        lineBraker(page,objs)

    return objs

def getInsertPosition(dst,ll):

    for n,d in enumerate(dst):
        if len(d[0])<ll:
            return n
    return n+1

def encyc(ds,ind):
    ret=[]
    for d in ds:
        ret.append(ind.index(d)+1)
    return ret

def mySort(src):
    dst=[];lm=0
    for r in src:
        ll=len(src[r][0])
        if len(dst)<1:
            dst.append([src[r][0],[r]+src[r][1]])
        else:
            nn=getInsertPosition(dst,ll)
            dst.insert(getInsertPosition(dst,ll),[src[r][0],[r]+src[r][1]])

    for d in dst:
#        print(len(d[0]),d)
        mm=max(d[0])
        if lm<mm:
                lm=mm

    index=[]
    for i in range(lm+1):
        index.append(-1)

    n=0
    for ds in dst:
#        print("ds",ds)
        for d in reversed(ds[0]):
            if index[d]==-1:
                index[d]=n;n+=1

    for ds in dst:
        for n,d in enumerate(ds[0]):
            ds[0][n]=index[d]

#    for d in dst:
#        print(len(d[0]),d)

    index=[]
    for ds in dst:
        for d in ds[1]:
            index.append(d)

#    for d in dst:
#        print(len(d[0]),d[0],encyc(d[1],index))

    return dst,index
## sort

def getNewList(newInfo):
    ret={}
    for s in newInfo:
        f=False;t=newInfo[s]
        for r in ret:
            if ret[r][0]==t:
                f=r
        if not f:
            ret[s]=t,[]
        else:
            ret[f][1].append(s)

    r2,idx=mySort(ret)

#    for r in r2:
#        print(len(r[0]),r[0],encyc(r[1],idx))
    return r2,idx

def getLink(smiles,index):
    ret=[];st=len(index)

    for s in reversed(smiles):
        if not s[-1] in index:
            index[s[-1]]=st;st+=1
        ret.append(index[s[-1]])

    s=smiles[0][0]
    if not s in index:
        index[s]=st;st+=1
    ret.append(index[s])

    ret.reverse()

    return ret

