import sys,os,time
import math
import subprocess
from rdkit import Chem
from rdkit import DataStructs
from rdkit.Chem.Fingerprints import FingerprintMols
from rdkit.Chem.Scaffolds import MurckoScaffold
from rdkit.Chem import rdFMCS
from rdkit.Chem import Draw, rdMolDescriptors
from rdkit.Chem.Draw import rdMolDraw2D
from rdkit.Chem import Descriptors
from rdkit.Chem import rdqueries
from IPython.display import SVG

from reportlab.graphics import renderPDF
from reportlab.pdfgen import canvas
from reportlab.pdfbase.pdfmetrics import stringWidth
from reportlab.lib.pagesizes import A4, portrait, landscape
from svglib.svglib import svg2rlg
from reportlab.lib.units import mm
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import shutil
import Levenshtein
import sqlite3
from dbTools import *
import re

import openpyxl
import pandas as pd
import unicodedata

eflag=False

def logout(mes):
    with open ("/var/www/html/public/images/report/step2.txt","a") as fp:
        if type(mes)==bool:
            if mes==False:
                fp.write("False\n")
            else:
                fp.write("True\n")
        if type(mes)==str:
            fp.write(mes+"\n")

def testForSvg(statAll):
    print(statAll[69])
    n=1
    fns=[];dst="/var/www/html/public/images/tmp/"
    for s in statAll[69][4]:
        if s!="":
            print(n,s)
            fn="now"+str(n)+".svg"
            smile2svg(s,fn)
            fns.append(fn)
            n+=1

    com=["montage","-tile",str(n-1)+"x1"]
    com=com+fns
    com.append("output.png")

    print(com)
    subprocess.Popen(com,cwd=dst)

def getPrm(cur,pids,sout):
    prm={};prmT={}
    for z in sout:
        prmT[z]={}
    for pid in pids:
        ret=cur.execute(f"""select * from searchList where id='{pid}';""")
        ret=ret.fetchall()
        if len(ret)<1:
            continue
        prm[pid]=ret[0]

    for k,v in prm.items():
        for z,x in zip(sout,v):
            prmT[z][k]=x

    return prm,prmT

def easyGetReaction(pid,cur):
    index={};newInfo={}
    sTable=f"searchTable{pid}"
    sql=f'select * from "{sTable}";'
    ans=cur.execute(sql).fetchall()
    for xx in reversed(ans):
        if xx[2]=="":
            continue
        newInfo[int(xx[0])]=getLink(strToSmiles(xx[2]),index)

    newList,idx=getNewList(newInfo)
#    print("-------------->",newList)
#    print("-------------->",idx,len(idx))
    hint=getReactionSummary(newList,idx)
#    print("\n\n\n");
#    print(hint)
    return hint

def checkIt(ss,ps):

    for s in ss:
        if not s in ps:
            return False
    return True

def inIt(h0,h1):
#    print("----->",h0,h1)
    l=len(h0)

    for n,h in enumerate(h1):
        if l<h[0]:
            if checkIt(h0,h[1]):
                return n
    return -1

def dropSubset(hint):
    ret={};rm=[];h1=[]

    for n,h0 in enumerate(hint[0]):
        h1.append([len(h0[1]),h0[1]])

    for n,h0 in enumerate(h1):
        m=inIt(h0[1],h1)
        if m>-1:
            rm.append(n)
    return rm

def getMax(h1):
    ret=0
    for h in h1:
        mx=max(h[1])
        if ret<mx:
            ret=mx
    return ret 

def makeNextIndex(simple_route,routes):
    nextIndex={}
    for route in simple_route:
        nextIndex[route]=routes.index(route)+1
    return nextIndex

def getOpeningMaterial(smiles):
    s=smiles[0]

    mol=Chem.MolFromSmiles(s)
    core=MurckoScaffold.GetScaffoldForMol(mol)
    coreSmiles=Chem.MolToSmiles(core)

    if coreSmiles=="" and smiles[1]!=">":
        s=smiles[1]
        mol=Chem.MolFromSmiles(s)
        core=MurckoScaffold.GetScaffoldForMol(mol)
        coreSmiles=Chem.MolToSmiles(core)
    return s, mol

def getLevel(s,sims):
    flag=False;lv=0
    for p in sims:
        if flag:
            if s==p:
                return lv
            if p=='(':
                lv=lv+1
            elif p==')':
                lv=lv-1
        if s==p:
            if ss=="[":
                continue
            flag=True
        ss=p
    return 0

def findRing(sims):
    flag=True;r="";lv=0
    for s in sims:
        if s.isdigit():
            if flag:# starting point
                if ss=="[":
                    continue
                p=s;flag=False;l=getLevel(s,sims);r=ss+s;continue
            if s==p:# end point
#                print("-->",l,r);
                r=r+s;return r
        if not flag:
            if s=="(":
                lv=lv+1;continue
            elif s==")":
                lv=lv-1;continue
#            print(lv,l,s,lv<=l)
            if lv <= l and not s.isdigit():
               # print("hit ",s)
                r=r+s
        ss=s
    return r

def prevFindRing(sims):
    r="";flag=True;level=0;p=1
    for s in sims:
        if s==")" or s=="]":
            level=level+1;continue
        if s=="(" or s=="[":
            level=level-1;continue
        if level!=0:
            continue
        if s.isdigit():
            if flag:
                p=s;r=ss+s;flag=False
            else:
                if s==p:
                    r=r+s;return r
                else:
                    print("unacceptable",sims)
                    r=r+s
        else:
            if not flag:
                r=r+s
        ss=s
    return r

def numRing(stars):
    level=0;n=0
#    print("starts",stars)
    for s in stars:
        if s=="[":
            level+=1;continue
        if s=="]":
            level-=1;continue
        if level!=0:
            continue
        if s.isdigit():
            n+=1
    return int(n/2)

def dict2Str(dct):
    sss="";
    for d in dct:
        sss=sss+str(d)+":"+str(dct[d])+","
    return sss

def getBinder(sims):
    global eflag
    cnt=0;tmp=sims

    while numRing(tmp)>0:
        r=findRing(tmp)
#        print("find",cnt,r,tmp);cnt+=1
        p=Chem.MolFromSmiles(tmp)
        x=Chem.MolFromSmiles(r)
#        tmp=Chem.MolToSmiles(Chem.ReplaceCore(p,x,labelByIndex=True))
        try:
            tmp=Chem.MolToSmiles(Chem.ReplaceCore(p,x))
        except:
            if eflag==False:
                eflag=True
                with open("statDb.err","a") as fp:
                    for z in sout:
                        sss=dict2Str(prmT[z])
                        fp.write(z+" "+sss+"\n")
                    fp.write("error in getBinder\n")
                    fp.write("smiles:"+tmp+" core:"+r+"\n")
            tmp=""
    return tmp

def countRing(sims):

    binder=""

    x=re.findall(r'\d',sims)
#    print(x,sims)
    n=int(len(x)/2)
    if n>1:
        binder=getBinder(sims)
    return int(len(x)/2),binder

def outOfRing(mol):
    num=mol.GetNumAtoms()
    n=0
    for i in range(num):
        if not mol.GetAtomWithIdx(i).IsInRing():
            n=n+1
    return n

def getMolSequence(allData):
    mols=[];smis=[]
#
# ugly patch:
#
    if allData['smiles'][0][0]=='>':
        tmp=['[Pu]']
        allData['smiles'][0]=tmp+allData['smiles'][0]
#    print(allData['smiles'][0])

    smiles,mol=getOpeningMaterial(allData['smiles'][0])
    mols.append(mol)
    smis.append(smiles)
    for s in allData['connect']:
        mols.append(Chem.MolFromSmiles(s[0]))
        smis.append(s[0])
    return mols,smis

def list2str(ss):
    d=""
    for s in ss:
        d=d+"["+str(s)+"]"
    return d

def reArrange(hint):
    stat=[{},{}]
    i=0
    maxLv=0
    print(hint[0])
    for h in hint[0]:
        if len(h[0])>maxLv:
            maxLv=len(h[0])

    if maxLv>5:
        maxLv=5
    for h in hint[0]:
        for lv in range(maxLv):
            lvs="row"+str(lv+1).zfill(2)
            if lv<len(h[0]):
                llv=lv
            else:
                llv=len(h[0])-1
            if not lvs in stat[i]:
                stat[i][lvs]=[h[0][llv]]
            else:
                stat[i][lvs].append(h[0][llv])
    i=1
    for lv in range(len(hint[0])):
        for h in hint[0][lv][0]:
            lvs="routeType"+str(lv+1).zfill(2)
            if not lvs in stat[i]:
                stat[i][lvs]=[h]
            else:
                stat[i][lvs].append(h)
    return stat

def analyzeCoreChanges(mols):
    step=0;p1="";p3=[];p4=[];p5=[]

    core=MurckoScaffold.GetScaffoldForMol(mols[0])
    coreSmiles=Chem.MolToSmiles(core)
    n=outOfRing(mols[0])
#    p3.append("o"+str(n))
    p3.append(n)
    m,bind=countRing(coreSmiles);p4.append(m);p5.append(bind)
    for mol in mols[1:]:
        nextCore=MurckoScaffold.GetScaffoldForMol(mol)
        nextCoreSmiles=Chem.MolToSmiles(nextCore)
        n=outOfRing(mol)
        if nextCoreSmiles!=coreSmiles:
            p1=p1+"x"
            coreSmiles=nextCoreSmiles;step=step+1
            p3.append(n)
#            p3.append("x"+str(n))
            m,bind=countRing(coreSmiles);p4.append(m);p5.append(bind)
        else:
            p1=p1+"o"
#            p3.append("o"+str(n))
            p3.append(n)
#    return [p1,len(mols)-1,p3,p4,p5,step]
    return [p1,len(mols)-1,p3,p4,p5]

def getAllIds(cur):
    ret=cur.execute(f"""select id from searchList;""").fetchall()
    if len(ret)>0:
        return ret

    print("maybe no records in this database")
    exit()

def getFirstIds(cur):
    pid=1
    while True:
        try:
            ret=cur.execute(f"""select * from searchList where id='{pid}';""")
        except:
            print("easy")

        if len(ret.fetchall())>0:
            return pid
        pid=pid+1
        if pid>50:
            print("maybe no records in this database")
            exit()

def xmerge(src):
    hand=False;cc=""
    if len(src)<1:
        return "Non"
    if len(src)==2:
        return src[0]+"-"+src[1]
    p=""
    for s in src:
        match s:
            case '(':
                hand=False;cc=cc+s
            case ')':
                hand=False;cc=cc+s
            case '=':
                hand=False;cc=cc+s
            case '.':
                hand=False;cc=cc+s
            case _:
#                if hand and "R" in s:
                if hand:
                    if 'R' in s:
                        cc=cc+"-"+s
                    elif 'R' in p:
                        cc=cc+"-"+s
                    else:
                        cc=cc+s
                else:
                    cc=cc+s
                hand=True
        p=s
    return cc

def xconv(src,i):
    match i:
        case 0:
            return src
        case 1:
            return str(src)
        case 2:
            return src
        case 3:
            return src

    cc="";depth=0;r="";rs=[]

    l=len(src)
    for n,s in enumerate(src):
        match s:
            case '[':
                depth=depth+1
                continue
            case ']':
                depth=depth-1
                if depth==0:
                    cc=cc+xmerge(rs)+","
                    rs=[]
                continue
            case '*':
                rs.append('R'+src[n-1])
            case _:
                if s.isdecimal():
                    continue
                rs.append(s)
        p=s
    dst=",".join(cc.split(",")[:-1])
#    print("\nsrc",src)
#    print("cc",dst)
    return dst

def sayThis(tp):
    match tp:
        case 1:
            print("core structure analyze")
            print("structure changed","number of Ring","number of atoms outside Rings")
            print("<-- route number (database), printed route number (in pdf)")
        case 2:
            print("route rearranging")
            print("list of routes","member of route","starting substance")
            print("<-- route number (database), printed route number (in pdf)")

def output(fp,com):

    if fp==True:
        print(com,end="")
    else:
        fp.write(com)

def addTable(page,data):
    page.setFont(font_size=3)
    sld=page.slide
    page.drawStringR(5,5,"this is ...xx");
    width=getWidth(data)
    print(data)
    df=pd.DataFrame(data)
    height=len(data)
    print(width[-1])
    if width[-1]>100:
        f=26.0/float(width[-1])
        add_table(sld,df,1,1,1+int(width[-1]*f),1+height,int(50*f),width)
    else:
        add_table(sld,df,2,2,2+width[-1]*2/6,2+height,18,width)

def xlen(sss):
    if 'IamAnImage' in sss:
        return 5
    i=0
    for s in sss:
        t=unicodedata.east_asian_width(s)
        if t=='H' or t=='Na' or t=='N':
            i=i+1
        else:
            i=i+2
    return i

def getWidth(data):

#    adjust..
    w=[0]*len(data)
    for n,k in enumerate(data):
        l=xlen(k)+2
        if w[n]<l:
            w[n]=l
        for d in data[k]:
            l=xlen(str(d))+2
            if w[n]<l:
                w[n]=l
    wx=0
    for i in w:
        wx=wx+i
    w.append(wx)
    return w

def setId(d,ys,k):
    d[k]=[]
    for i in ys:
        d[k].append(i)

def setDDD(d,ys,k,n):
#komai
    d[k]=[]
    for y in ys:
        if n==1:
            d[k].append("{:.2f}".format(ys[y][n]*100.0))
        elif n==2:
            d[k].append("{:.2f}".format(ys[y][n]))
        elif n==3:
            d[k].append(int(ys[y][n]))
        else:
            d[k].append(str(ys[y][n]))

def setPic(d,p,k,s,sl):
    d[k]=[]
    for i in sl:
        smiles=p[s][i]
        dr="/var/www/html/public/images/tmp/"
        if not os.path.exists(dr):
            os.makedirs(dr)
        fn=dr+str(i)+".svg"

        mol=Chem.MolFromSmiles(smiles)
        view = rdMolDraw2D.MolDraw2DSVG(200,200)

        option=view.drawOptions()
        option.circleAtoms=False
        option.continuousHightlight=False

        tm = rdMolDraw2D.PrepareMolForDrawing(mol)
        view.DrawMolecule(tm)
        view.FinishDrawing()

        svg=view.GetDrawingText()

        print(smiles,fn)

        with open(fn,'w') as f:
            f.write(svg)
        cropSvg(fn)
        gn=fn.split(".svg")[0]+".jpg"
        subprocess.run(['/usr/bin/convert','-density','360',fn,gn])
        d[k].append('IamAnImage:'+gn)

def setDat(d,p,k,s,sl):
    d[k]=[]
    for i in sl:
        if s=='factors' or s=='options':
            r="".join(p[s][i].split("'"))
            d[k].append(r)
        else:
            d[k].append(p[s][i])

def getCellPosition(table,row,col):
    table_gf = table._graphic_frame
    x = table_gf._element.xfrm.off.x
    y = table_gf._element.xfrm.off.y

    cell_left = x
    cell_top = y
    cell_height = table.rows[0].height
    cell_width = table.columns[0].width

    for r in range(row):
        cell_height = table.rows[r].height
        cell_top += cell_height*1.55

    for c in range(col):
        cell_width = table.columns[c].width
        cell_left += cell_width

    cell_height = table.rows[row].height
    cell_width = table.columns[col].width

    return cell_left, cell_top, cell_width, cell_height

#komai
def countTheAtom(mol,C):
    cnt=0
    for atom in mol.GetAtoms():
        if atom.GetSymbol()==C:
            cnt=cnt+1
    return cnt

def isOrganicHalogen(smiles):
    tg=['F','Cl','Br','I','At']
    r=[];c=0

    mol=Chem.MolFromSmiles(smiles)
    for atom in mol.GetAtoms():
        s=atom.GetSymbol()
        if s=='C':
            c=c+1
        if s in tg:
            if not s in r:
                r.append(s)
    if c>0 and len(r)>0:
        return True,r
    return False,r

def fullCouplingList(theRoute,allData):
    hit={};f=[]
    for route in theRoute:
        flag=True;c0=0;n0=0
        for m,smiles in enumerate(allData[route]['smiles']):
#            print("--->",m,smiles)
            l=len(smiles)-1
#            if l<3:
            if l!=3:
                continue
            for i in range(l-1):
                now=smiles[i]
                check,halogens=isOrganicHalogen(now)
#                print("isOrganic",f,now,i,l)
                if check:
#                    if chkCoupling(f,now,smiles[l],smiles):
                    f,cc=chkCoupling(halogens,now,smiles[l],smiles)
                    if f:
                        if flag:
                            print("route <-------------------------",route)
                            flag=False
                        if not route in hit:
                            hit[route]=[]
                        hit[route].append([m,now,smiles[1-i],smiles[3]])
#                        hit[route].append([m,now,cc[0]+"-"+cc[1]+"+"+cc[2]+"-"+cc[3]+"->"+smiles[3]])

    return {'crossCoupling':hit}


def isOrganicMetal(mol):
    ret=[]

    metal_symbols = {
            'Mg','Zn','Sn','B','Si','Pd'
            }
    xmetal_symbols = {
        'Li', 'Be', 'Na', 'Mg', 'Al', 'K', 'Ca', 'Sc', 'Ti', 'V', 'Cr', 'Mn',
        'Fe', 'Co', 'Ni', 'Cu', 'Zn', 'Ga', 'Ge', 'As', 'Se', 'Br', 'Kr',
        'Rb', 'Sr', 'Y', 'Zr', 'Nb', 'Mo', 'Tc', 'Ru', 'Rh', 'Pd', 'Ag', 'Cd',
        'In', 'Sn', 'Sb', 'Te', 'I', 'Xe', 'Cs', 'Ba', 'La', 'Ce', 'Pr', 'Nd',
        'Pm', 'Sm', 'Eu', 'Gd', 'Tb', 'Dy', 'Ho', 'Er', 'Tm', 'Yb', 'Lu', 'Hf',
        'Ta', 'W', 'Re', 'Os', 'Ir', 'Pt', 'Au', 'Hg', 'Tl', 'Pb', 'Bi', 'Po',
        'At', 'Rn', 'Fr', 'Ra', 'Ac', 'Th', 'Pa', 'U', 'Np', 'Pu', 'Am', 'Cm',
        'Bk', 'Cf', 'Es', 'Fm', 'Md', 'No', 'Lr', 'Rf', 'Db', 'Sg', 'Bh', 'Hs',
        'Mt', 'Ds', 'Rg', 'Cn', 'Nh', 'Fl', 'Mc', 'Lv', 'Ts', 'Og'
    }

    for atom in mol.GetAtoms():
        if atom.GetSymbol() in metal_symbols:
            ret.append(atom.GetSymbol())
    return ret

def chkByHalogenNum(smiles,f):
    flag=False;n=0
    for s in smiles:
        if s=='>':
            flag=True
        else:
            m=countTheAtom(Chem.MolFromSmiles(s),f)
            if flag:
                n=n-m
            else:
                n=n+m
    if n>0:
        return True
    return False

def chkCoupling(halogens,material,product,smiles):
#    print(smiles)
    mol1=Chem.MolFromSmiles(material)
    mol2=Chem.MolFromSmiles(product)
    flag=False
    for f in halogens:
        M=False;R2=False
        R1=Chem.DeleteSubstructs(mol1,Chem.MolFromSmiles(f))
        if not chkByHalogenNum(smiles,f):
            continue
        if mol2.HasSubstructMatch(R1):
            if material==smiles[0]:
                mp=smiles[1]
                material_pair=Chem.MolFromSmiles(smiles[1])
            else:
                mp=smiles[0]
                material_pair=Chem.MolFromSmiles(smiles[0])

#            R3=Chem.DeleteSubstructs(mol2,R1)
#            R3s=Chem.MolToSmiles(R3)
            mols=[mol2,material_pair]
            res=rdFMCS.FindMCS(mols)
            R2=Chem.MolFromSmarts(res.smartsString)

            mols=[mol2,mol1]
            res=rdFMCS.FindMCS(mols)
            R1a=Chem.MolFromSmarts(res.smartsString)

            M=Chem.DeleteSubstructs(material_pair,R2)
#            print(material,product,R3s)
#            print("R2",Chem.MolToSmiles(R2),"R3",Chem.MolToSmiles(R3))
#            print(R2s,R3s)

            if isOrganicMetal(M):
                c2=True
            else:
                c2=False

            if c2:
                flag=True
                print(smiles)
                print("material",material,f)
                print("pair    ",mp)
                print("product ",product)
                print("R1      ",Chem.MolToSmiles(R1))
                print("R1a     ",Chem.MolToSmiles(R1a))
                print("R2      ",Chem.MolToSmiles(R2))
                print("M       ",Chem.MolToSmiles(M))
#        else:
#            print("not:",Chem.MolToSmiles(mol2),Chem.MolToSmiles(R1))
    return flag, [f,R1,M,R2]

def couplingList(theRoute,allData):
    hit=[];f=[]
    for route in theRoute:
        flag=True;c0=0;n0=0
        for m,smiles in enumerate(allData[route]['smiles']):
            l=len(smiles)-1
            if l!=3:
                continue
            for i in range(l-1):
                now=smiles[i]
                check,halogens=isOrganicHalogen(now)
                if check:
                    f,_=chkCoupling(halogens,now,smiles[l],smiles)
                    if f:
                        if flag:
                            flag=False
                        if not route in hit:
                            hit.append(route)
    return {'crossCoupling':hit}

def makeShowList(ys):
    sl=[]
    for k in ys:
        sl.append(k)
    return sl

def add_table(sld, df, left, top, width, height, font_size, width_index):
    '''
    args:
        slide[slide]: Slide object
        df[DataFrame] : Display data
        left[int]: Position from the left end
        top[int] : Position from top
        width[int]: Width of object
        height[int]: Height of object
        font_size[int]: Font size
    return:
        None
    '''
    dr="/var/www/html/public/images/tmp/"

    shapes=sld.shapes
    column_names = df.columns.tolist()
    index_names = df.index.tolist()

    col_num = len(column_names)
    row_num = len(index_names) + 1

    table = shapes.add_table(
         row_num,
         col_num,
         Cm(left),
         Cm(top),
         Cm(width),
         Cm(height)
    ).table

    # Columnの値を挿入
    for i in range(col_num):
        table.columns[i].width = Cm(width*width_index[i]/width_index[-1])
        table.cell(0, i).text = column_names[i]
        table.cell(0, i).text_frame.paragraphs[0].font.size = Pt(font_size)
            # 中央揃え
        table.cell(0, i).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # セルの値を挿入
    for i in range(1, row_num):
        for j in range(col_num):
            if 'IamAnImage' in str(df.iloc[i-1,j]):
                fn=df.iloc[i-1,j].split(":")[1]
                with open(fn,"rb") as fp:
                    img=fp.read()
                x,y,w,h=getCellPosition(table,i,j)
                shapes.add_picture(fn,x,y,width=w,height=h)
            else:
                table.cell(i, j).text = str(df.iloc[i-1, j])
            table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(font_size)
            # 中央揃え
            if type(df.iloc[i-1,j])==str:
                table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            else:
                table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
#            table.cell(i, j).text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT


#########################################################
###################### main #############################
#########################################################
input_dir="./";db="sList.db"
com='-chem chemical_name, -ppt for power point -n [1,2,3] to specify routes -p proc_per_line -d input_path, -one_size : rendering at the same size, -product_only : categolize routes with respct to reaction product only, -show_subsets : show routes including subsets'
com+=" -summary 0/1/2"
ops={'-h':com}

ppt=False;skip=False;hit_id=False;stdout=True
ids=[];mode="normal";full=True;tp=1;label=False;fn=False
excel_file_name=False;items="1,2,3";

for n,op in enumerate(sys.argv):
    if hit_id:
        if op[0]!='-':
            ids.append(int(op))
        else:
            hit_id=False
    if skip:
        skip=False;continue
    match op:
        case '-h':
            print(sys.argv[0],ops['-h']);exit()
        case '-v':
            print(version);exit()
        case '-ppt':
            ppt=True
        case '-items':
            items=sys.argv[n+1]
            skip=True
        case '-type':
            tp=int(sys.argv[n+1])
            skip=True
        case '-simple':
            full=False
        case '-ids':
            hit_id=True
        case '-label':
            label=sys.argv[n+1].split(',')
        case '-database':
            db=sys.argv[n+1]
            skip=True
        case '-forEval':
            mode="forEval"
        case '-stdout':
            stdout=open(sys.argv[n+1],"w")
            skip=True
        case '-excel_file_name':
            excel_file_name=sys.argv[n+1]
        case '-d':
            input_dir=sys.argv[n+1]
            skip=True

conn=sqlite3.connect(db)
cur=conn.cursor()
error=False

if len(ids)<1:
    ids.append(getFirstIds(cur))
elif ids[0]==0:
    ret=getAllIds(cur)
    ids=[]
    for id in ret:
        ids.append(id[0])
    pid=ids[0]

print("ids",ids)

ret=cur.execute(f"""select * from searchList;""")
descr=cur.description

if tp!=4:
    page=openReport("test.txt",0.5)
    page.setOutput(stdout);

sout=[]
for d in descr:
    sout.append(d[0])

prm,prmT=getPrm(cur,ids,sout)

for z in sout:
    print(z,prmT[z])

hint={}
for id in ids:
    hint[id]=easyGetReaction(id,cur)

cutList={}
for id in ids:
    cutList[id]=dropSubset(hint[id])
#    print("hint\n")
#    print(hint[id],len(hint[id]))
#    print("cutList\n")
#    print(cutList[id])

#for h in hint[id]:
#    print("h",h,"\n",len(h),"\n")
#exit()
result={}
ys={}
for id in ids:
    ret=cur.execute(f"""select user,date from searchList where id='{id}';""")
    ret=ret.fetchall()

    uid=ret[0][0];date=ret[0][1];
    dst="/var/www/html/public/images/"+uid+"/"+date+"/route"

    result[id]={}
    m=getMax(hint[id][0])
#    print(len(hint[id][0]),m,len(hint[id][0])-len(cutList[id]))
    result[id]['routeNumber']=len(hint[id][0])
    result[id]['chemicalNumber']=m
    result[id]['fullRouteNumber']=len(hint[id][0])-len(cutList[id])

#    print(id,hint[id])

ox=40
fcl={}
##############################################################
for id in ids:
    sTable=f"searchTable{id}"
    ans=cur.execute(f"""select * from {sTable};""").fetchall()
    if len(ans[0][1].split(";"))<2:
        output(stdout,f"no route {id}\n")
        if tp==3:
            ys[id]=[0,0.0,0,0]
        print(f"no route {id}");continue

    allData={};routes=[];newInfo={};index={}
    for xx in ans:
        s={}
        if xx[1]=="":
            continue
        routes.append(int(xx[0]))
        s['connect']=strToConnect(xx[1])
        s['smiles']=strToSmiles(xx[2])
        s['info']=strToSmiles(xx[3])
        allData[int(xx[0])]=s

    for xx in reversed(ans):
        newInfo[int(xx[0])]=getLink(strToSmiles(xx[2]),index)

    newList,idx=getNewList(newInfo)

    simple_route=[]
    for nn in newList:
        simple_route.append(nn[1][0])

    full_route=[]
    for nn in newList:
        full_route=full_route+nn[1]

    if full:
        theRoute=full_route
    else:
        theRoute=simple_route

    hint=getReactionSummary(newList,idx)
    nextIndex=makeNextIndex(theRoute,hint[3])

    statAll={};
    if tp==3 or tp==4:
        p1=len(full_route)
        p2=float(p1)/float(prmT['loop'][id])
        p3=0
        allSmis=[]
        for route in theRoute:
            mols,smis=getMolSequence(allData[route])
            for s in allData[route]['smiles']:
                if s not in allSmis:
                    allSmis.append(s)
            p3=p3+len(mols)-1
        p3=float(p3)/float(p1)
        ys[id]=[p1,p2,p3,len(allSmis)]

    match tp:
#        case 4:
#        case 3:
#            print(id,result[id]['routeNumber'],result[id]['chemicalNumber'],result[id]['fullRouteNumber']);
        case 1:
            stat=[{},{},{},{},{}]
            ppp=0
            for route in theRoute:
                mols,smis=getMolSequence(allData[route])
#    [coreChange,coreCNum,ringNum,outOfRings]=analyzeCoreChanges(mols)
                r=analyzeCoreChanges(mols)
#                print(r[0],list2str(r[3]),list2str(r[2]),list2str(r[4]),"<--",route,nextIndex[route],r[1])

                statAll[route]=r
                for i in range(len(r)):
                    if type(r[i])==list:
                        st=list2str(r[i])
                    else:
                        st=r[i]

                    if not st in stat[i]:
                        stat[i][st]=[route]
                    else:
                        stat[i][st].append(route)
        case 2:
            stat=reArrange(hint)
            stat.append(couplingList(theRoute,allData))
            print(stat[-1])
            ppp=1
        case 5:
            fcl[id]=fullCouplingList(theRoute,allData)

    if tp==3 or tp==4:
        continue
    if tp==5:
        continue

    output(stdout,f"start-forEval for {id}\n")
    for i in range(len(stat)):
#        print(stat[i])
        output(stdout,"key"+str(i+1)+"\n")
        for key in sorted(stat[i]):
            output(stdout,"##"+xconv(key,i)+"#")
            for route in sorted(stat[i][key]):
                output(stdout,str(route)+",")
        output(stdout,"\n")
    output(stdout,"end-forEval\n")
#
#  for drawImages
#
    output(stdout,f"start-forDraw for {id}\n")
    for key in sorted(stat[ppp]):
        for route in stat[ppp][key]:
            ddd=dst+"/"+str(route)+"/"
#            print("route"+str(route)+"\npath:"+ddd.split("/var/www/html/public")[1],end=";")
            output(stdout,"route"+str(route)+"\npath:"+ddd.split("/var/www/html/public")[1]+";")
            xx=ox
            connect=allData[route]['connect']
            smiles=allData[route]['smiles']
            info=allData[route]['info']
            fc,image_size=makeSvg(smiles,ddd)
            objs=makeObject(page,connect,image_size)
            for obj in objs:
                for item in obj:
                    x0=item['pos'][0][0];y0=item['pos'][0][1]
                    if item['type']==1:
                        page.drawImage(ddd+"/"+item['svg'][0],x0,y0)
                    else:
                        x1=item['pos'][1][0];y1=item['pos'][1][1]
                        xx=x0;

                        for n,fn in enumerate(item['svg']):
                            xx=xx+page.mg
                            r=image_size[fn]
                            page.drawImage(ddd+fn,xx,y0+page.mg,cat=True)
                            xx+=r[0]
                        page.drawArrow(arrow(x0,y0,x1,y1,page.scale))
            page.stroke()
            output(stdout,"\n\n");
    output(stdout,"end-forDraw\n")

if tp==5:
    id_head=True
    print(fcl)
    output(stdout,f"start-forCrossCoupling\n")
    output(stdout,'{"title":"Reaction with CrossCoupling",\n"body":\n{')
    for xid in fcl:
        cc=fcl[xid]['crossCoupling']
        if len(cc)<1:
            continue
        if id_head!=True:
            output(stdout,',\n')
        output(stdout,'"id'+str(xid)+'":{\n')
        id_head=False
#        print("-->",fcl[xid])
        headFlag1=True
        if headFlag1!=True:
            output(stdout,",")
        else:
            headFlag1=False
        headFlag2=True
        for pop in cc:
            if headFlag2!=True:
                output(stdout,",\n")
            else:
                headFlag2=False
#            print(pop)
            output(stdout,'"route'+str(pop)+'":{');
            headFlag3=True
            for c in cc[pop]:
                if headFlag3!=True:
                    output(stdout,",\n")
                else:
                    headFlag3=False
#                output(stdout,'"step'+str(c[0]+1)+'":{"smiles":"'+c[1]+'"}');
#                output(stdout,'"step'+str(c[0]+1)+'":"'+c[1]+'"');
                output(stdout,'"step'+str(c[0]+1)+'":"'+c[2]+'"');
#                output(stdout,'"step'+str(c[0]+1)+'":"'+c[1]+'+'+c[2]+'->'+c[3]+'"');
            output(stdout,"}")
        output(stdout,"}")
    output(stdout,"}}\n")
    output(stdout,f"end-forCrossCoupling\n")

if tp==3:
    title=["探索数","探索率","平均ステップ数","出現物質数"]
    lll=[]
    for l in label:
        lll.append(",".join(l.split("__")))

    for y in ys:
        if excel_file_name==False:
            excel_file_name="test.xlsx"
    saveExcel(ys,lll,title,input_dir+"/"+excel_file_name)

if stdout!=True:
    stdout.close()
    stdout=True

if tp==4:
    test_data = {
    'Name': ['Alice', 'Bob', 'Charlie'],
    'Age': ['25', '30', '35'],
    'Occupation': ['Engineer', 'Designer', 'Doctor']
    }
    data={}
    logout("items:"+items)
    title=["探索数","探索率","平均ステップ数","出現物質数"]
    title=["探索数","探索率","ステップ数","物質数"]

    print("<----------------------------------")
    print(prmT)
    print("<----------------------------------")
    print(ys)
    print("<----------------------------------")
    sl=makeShowList(ys)
    items=items.split(',')
    logout('0' in items)
    if '0' in items:
        setId(data,ys,'id',sl)
    if '1' in items:
        setDat(data,prmT,'物質名','substance',sl)
    if '2' in items:
        setDat(data,prmT,'化学式','smiles',sl)
    if '3' in items:
        setPic(data,prmT,'構造式','smiles',sl)
    if '4' in items:
        setDat(data,prmT,'要求数','loop',sl)
    if '5' in items:
        setDDD(data,ys,'探索数',0)
    if '6' in items:
        setDDD(data,ys,'探索率',1)
    if '7' in items:
        setDDD(data,ys,'合成ステップ数',2)
    if '8' in items:
        setDDD(data,ys,'出現物質数',3)
    if '9' in items:
        setDat(data,prmT,'重み','factors',sl)
    if '10' in items:
        setDat(data,prmT,'条件','options',sl)
#    $item=array("id","物質名","SMILES","構造式(画像)","要求数","獲得数","探索率","合成ステップ数","出現物質数","探索経路の重み","探索条件");

    page=openReport(input_dir+"/"+uid+"Table.pptx",0.5,title="table")

    addTable(page,data)
    page.close()
#
#        print(core)
#        print(lastP)
#        print("info",s['info'])
